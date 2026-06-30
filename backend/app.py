# -*- coding: utf-8 -*-
"""
PDFMaster - 全能PDF处理软件 后端服务
功能覆盖：查看、编辑、注释、转换、压缩、合并拆分、OCR、水印、签名、加密、打印
"""

import os
import io
import json
import uuid
import base64
import shutil
import zipfile
from pathlib import Path

from flask import Flask, request, jsonify, send_file, render_template
from flask_cors import CORS

import fitz  # PyMuPDF
from PIL import Image
import pikepdf

# 配置
BASE_DIR = Path(__file__).resolve().parent.parent
UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "output"
UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

app = Flask(__name__, static_folder=str(BASE_DIR / "frontend"), static_url_path="")
CORS(app)


# ============================================================
# 工具函数
# ============================================================

def gen_id():
    return uuid.uuid4().hex[:12]


def save_upload(file):
    """保存上传文件"""
    fid = gen_id()
    ext = Path(file.filename).suffix.lower()
    fpath = UPLOAD_DIR / f"{fid}{ext}"
    file.save(str(fpath))
    return str(fpath), fid


def json_ok(data=None, msg="success"):
    return jsonify({"code": 0, "msg": msg, "data": data})


def json_err(msg, code=1):
    return jsonify({"code": code, "msg": msg, "data": None})


# ============================================================
# 首页
# ============================================================

@app.route("/")
def index():
    return send_file(str(BASE_DIR / "frontend" / "index.html"))


# ============================================================
# 1. PDF 查看/阅读
# ============================================================

@app.route("/api/pdf/info", methods=["POST"])
def pdf_info():
    """获取PDF基本信息"""
    f = request.files.get("file")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        toc = doc.get_toc()
        info = {
            "fid": fid,
            "filename": f.filename,
            "pages": doc.page_count,
            "size": os.path.getsize(fpath),
            "metadata": doc.metadata,
            "toc": toc,
        }
        # 获取每页尺寸
        pages_info = []
        for i in range(doc.page_count):
            page = doc[i]
            rect = page.rect
            pages_info.append({
                "page": i + 1,
                "width": rect.width,
                "height": rect.height,
            })
        info["pages_info"] = pages_info
        doc.close()
        return json_ok(info)
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/page/<fid>/<int:page_num>")
def pdf_page_image(fid, page_num):
    """获取PDF单页渲染图片"""
    fpath = _find_upload(fid)
    if not fpath:
        return json_err("文件不存在")
    zoom = float(request.args.get("zoom", 1.5))
    try:
        doc = fitz.open(fpath)
        if page_num < 0 or page_num >= doc.page_count:
            doc.close()
            return json_err("页码超出范围")
        page = doc[page_num]
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        img_data = pix.tobytes("png")
        doc.close()
        return send_file(io.BytesIO(img_data), mimetype="image/png")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/text/<fid>")
def pdf_extract_text(fid):
    """提取PDF文本"""
    fpath = _find_upload(fid)
    if not fpath:
        return json_err("文件不存在")
    try:
        doc = fitz.open(fpath)
        pages_text = []
        for i in range(doc.page_count):
            page = doc[i]
            text = page.get_text("text")
            pages_text.append({"page": i + 1, "text": text})
        doc.close()
        return json_ok(pages_text)
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 2. PDF 编辑
# ============================================================

@app.route("/api/pdf/edit/text", methods=["POST"])
def pdf_edit_text():
    """编辑PDF文本(添加/删除)"""
    f = request.files.get("file")
    edits_str = request.form.get("edits", "[]")
    if not f:
        return json_err("请上传文件")
    edits = json.loads(edits_str)
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        for edit in edits:
            page_num = edit["page"] - 1
            action = edit.get("action", "add")
            if action == "add":
                page = doc[page_num]
                x, y = edit.get("x", 100), edit.get("y", 100)
                text = edit.get("text", "")
                fontsize = edit.get("fontsize", 12)
                color = _parse_color(edit.get("color", "#000000"))
                fontname = edit.get("fontname", "helv")
                page.insert_text((x, y), text, fontname=fontname, fontsize=fontsize, color=color)
            elif action == "delete":
                page = doc[page_num]
                # 用白色方块覆盖文字
                x, y, w, h = edit.get("x", 0), edit.get("y", 0), edit.get("w", 100), edit.get("h", 20)
                page.draw_rect(fitz.Rect(x, y, x + w, y + h), color=(1, 1, 1), fill=(1, 1, 1))
        out_path = _output_path(fid)
        doc.save(out_path)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name=f"edited_{f.filename}")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/edit/image", methods=["POST"])
def pdf_insert_image():
    """PDF中插入图片"""
    f = request.files.get("file")
    img_file = request.files.get("image")
    if not f or not img_file:
        return json_err("请上传PDF和图片")
    page_num = int(request.form.get("page", 1))
    x = float(request.form.get("x", 100))
    y = float(request.form.get("y", 100))
    w = float(request.form.get("w", 100))
    h = float(request.form.get("h", 100))

    fpath, fid = save_upload(f)
    img_bytes = img_file.read()

    try:
        doc = fitz.open(fpath)
        page = doc[page_num - 1]
        img_rect = fitz.Rect(x, y, x + w, y + h)
        page.insert_image(img_rect, stream=img_bytes)
        out_path = _output_path(fid)
        doc.save(out_path)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name=f"img_{f.filename}")
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 3. PDF 注释/标注
# ============================================================

@app.route("/api/pdf/annotate", methods=["POST"])
def pdf_annotate():
    """添加注释到PDF"""
    f = request.files.get("file")
    annotations_str = request.form.get("annotations", "[]")
    if not f:
        return json_err("请上传文件")
    annotations = json.loads(annotations_str)
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        for ann in annotations:
            page_num = ann["page"] - 1
            ann_type = ann.get("type", "highlight")
            page = doc[page_num]
            rect = fitz.Rect(ann["x1"], ann["y1"], ann["x2"], ann["y2"])

            if ann_type == "highlight":
                annot = page.add_highlight_annot(rect)
            elif ann_type == "underline":
                annot = page.add_underline_annot(rect)
            elif ann_type == "strikeout":
                annot = page.add_strikeout_annot(rect)
            elif ann_type == "squiggly":
                annot = page.add_squiggly_annot(rect)
            elif ann_type == "text":
                annot = page.add_text_annot(rect.tl, ann.get("content", ""))
            elif ann_type == "free_text":
                annot = page.add_freetext_annot(rect, ann.get("content", ""),
                    fontsize=ann.get("fontsize", 12))
            elif ann_type == "rect":
                color = _parse_color(ann.get("color", "#FF0000"))
                annot = page.add_rect_annot(rect)
            elif ann_type == "sticky_note":
                annot = page.add_text_annot(rect.tl, ann.get("content", ""))

            if annot and ann.get("color"):
                try:
                    annot.set_colors(stroke=_parse_color_rgb(ann["color"]))
                except:
                    pass
        out_path = _output_path(fid)
        doc.save(out_path)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name=f"annotated_{f.filename}")
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 4. PDF 转换
# ============================================================

@app.route("/api/pdf/convert/to-word", methods=["POST"])
def pdf_to_word():
    """PDF转Word"""
    f = request.files.get("file")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        from docx import Document
        doc_pdf = fitz.open(fpath)
        doc_word = Document()
        for i in range(doc_pdf.page_count):
            page = doc_pdf[i]
            text = page.get_text("text")
            if text.strip():
                doc_word.add_paragraph(text)
            # 提取图片
            images = page.get_images()
            for img in images:
                try:
                    xref = img[0]
                    base_image = doc_pdf.extract_image(xref)
                    img_bytes = base_image["image"]
                    img_ext = base_image["ext"]
                    img_stream = io.BytesIO(img_bytes)
                    doc_word.add_picture(img_stream)
                except:
                    pass
        out_path = _output_path(fid, ".docx")
        doc_word.save(out_path)
        doc_pdf.close()
        return send_file(out_path, as_attachment=True, download_name=f"{Path(f.filename).stem}.docx")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/convert/to-excel", methods=["POST"])
def pdf_to_excel():
    """PDF转Excel - 提取表格"""
    f = request.files.get("file")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        from openpyxl import Workbook
        doc = fitz.open(fpath)
        wb = Workbook()
        ws = wb.active
        row_idx = 1
        for i in range(doc.page_count):
            page = doc[i]
            tables = page.find_tables()
            if tables.tables:
                for table in tables.tables:
                    for row in table.extract():
                        for col_idx, cell in enumerate(row, 1):
                            ws.cell(row=row_idx, column=col_idx, value=str(cell) if cell else "")
                        row_idx += 1
                    row_idx += 1
            else:
                text = page.get_text("text")
                for line in text.split("\n"):
                    if line.strip():
                        ws.cell(row=row_idx, column=1, value=line)
                        row_idx += 1
        out_path = _output_path(fid, ".xlsx")
        wb.save(out_path)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name=f"{Path(f.filename).stem}.xlsx")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/convert/to-ppt", methods=["POST"])
def pdf_to_ppt():
    """PDF转PPT"""
    f = request.files.get("file")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        from pptx import Presentation
        from pptx.util import Inches
        doc = fitz.open(fpath)
        prs = Presentation()
        for i in range(doc.page_count):
            page = doc[i]
            pix = page.get_pixmap(dpi=150)
            img_bytes = pix.tobytes("png")
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, Inches(0), Inches(0),
                                     width=prs.slide_width, height=prs.slide_height)
        out_path = _output_path(fid, ".pptx")
        prs.save(out_path)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name=f"{Path(f.filename).stem}.pptx")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/convert/to-image", methods=["POST"])
def pdf_to_image():
    """PDF转图片"""
    f = request.files.get("file")
    fmt = request.form.get("format", "png")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        images = []
        for i in range(doc.page_count):
            page = doc[i]
            pix = page.get_pixmap(dpi=200)
            img_bytes = pix.tobytes(fmt)
            images.append(img_bytes)
        doc.close()

        if len(images) == 1:
            return send_file(io.BytesIO(images[0]), mimetype=f"image/{fmt}",
                           download_name=f"{Path(f.filename).stem}.{fmt}")
        else:
            # 打包为ZIP
            zip_buf = io.BytesIO()
            with zipfile.ZipFile(zip_buf, "w") as zf:
                for i, img_data in enumerate(images):
                    zf.writestr(f"page_{i+1}.{fmt}", img_data)
            zip_buf.seek(0)
            return send_file(zip_buf, mimetype="application/zip",
                           download_name=f"{Path(f.filename).stem}_images.zip")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/convert/to-html", methods=["POST"])
def pdf_to_html():
    """PDF转HTML"""
    f = request.files.get("file")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        html_parts = ['<html><head><meta charset="utf-8"><style>',
                      'body{font-family:SimSun,sans-serif;max-width:800px;margin:auto;padding:20px}',
                      '.page{border-bottom:1px dashed #ccc;padding:10px 0;page-break-after:always}',
                      '</style></head><body>']
        for i in range(doc.page_count):
            page = doc[i]
            html_parts.append(f'<div class="page"><h3>第{i+1}页</h3>')
            html_parts.append(page.get_text("html"))
            html_parts.append('</div>')
        html_parts.append('</body></html>')
        doc.close()
        out_path = _output_path(fid, ".html")
        with open(out_path, "w", encoding="utf-8") as fh:
            fh.write("\n".join(html_parts))
        return send_file(out_path, as_attachment=True, download_name=f"{Path(f.filename).stem}.html")
    except Exception as e:
        return json_err(str(e))


# ---- X-to-PDF 专用端点 ----
@app.route("/api/pdf/convert/to-pdf/word", methods=["POST"])
def word_to_pdf():
    """Word (.docx) 转 PDF"""
    f = request.files.get("file")
    if not f: return json_err("请上传Word文件")
    fpath, fid = save_upload(f)
    try:
        from docx import Document as DocxDoc
        docx = DocxDoc(fpath)
        doc = fitz.open()
        A4_W, A4_H = 595, 842
        margin = 60
        y = margin
        page = doc.new_page()

        def new_page_if_needed(y_pos, needed=24):
            nonlocal page
            if y_pos + needed > A4_H - margin:
                page = doc.new_page()
                return margin
            return y_pos

        for para in docx.paragraphs:
            text = para.text.strip()
            if not text:
                y += 12
                y = new_page_if_needed(y)
                continue
            # 标题样式
            fs = 11
            if para.style.name.startswith("Heading") or para.style.name.startswith("标题"):
                level = 1
                try: level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1
                except: pass
                fs = {1: 20, 2: 16, 3: 14}.get(level, 14)
                y = new_page_if_needed(y, fs + 12)
                page.insert_text((margin, y), text, fontsize=fs, fontname="hebo", fontfile=None)
                y += fs + 12
            else:
                y = new_page_if_needed(y, fs + 6)
                page.insert_text((margin, y), text, fontsize=fs, fontname="hebo", fontfile=None)
                y += fs + 6

        # 处理表格
        for table in docx.tables:
            y = new_page_if_needed(y, 50)
            for row_idx, row in enumerate(table.rows):
                cols = [cell.text.strip()[:40] for cell in row.cells]
                x = margin
                col_w = (A4_W - 2 * margin) / max(len(cols), 1)
                for cell_text in cols:
                    page.draw_rect(fitz.Rect(x - 2, y - 2, x + col_w - 2, y + 20),
                                   color=(0.6, 0.6, 0.6), width=0.5)
                    page.insert_text((x, y + 12), cell_text, fontsize=9, fontname="hebo", fontfile=None)
                    x += col_w
                y += 24
                y = new_page_if_needed(y, 24)

        out_path = _output_path(fid)
        doc.save(out_path, garbage=4, deflate=True)
        doc.close()
        return send_file(out_path, as_attachment=True,
                        download_name=f"{Path(f.filename).stem}.pdf")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/convert/to-pdf/excel", methods=["POST"])
def excel_to_pdf():
    """Excel (.xlsx) 转 PDF —— 每个sheet为一页或多页表格"""
    f = request.files.get("file")
    if not f: return json_err("请上传Excel文件")
    fpath, fid = save_upload(f)
    try:
        from openpyxl import load_workbook
        wb = load_workbook(fpath, data_only=True)
        doc = fitz.open()
        A4_W, A4_H = 595, 842
        margin = 36

        for ws in wb.worksheets:
            rows = [[cell.value for cell in row] for row in ws.iter_rows(
                max_col=min(ws.max_column or 1, 10))]
            if not rows: continue
            # 过滤全空行
            rows = [r for r in rows if any(c is not None for c in r)]
            if not rows: continue

            ncols = max(len(r) for r in rows)
            max_text_w = max(80, (A4_W - 2 * margin) / ncols)
            row_h = 20
            y = margin + 30

            page = doc.new_page()
            page.insert_text((margin, margin), ws.title, fontsize=14, fontname="hebo", fontfile=None)

            for ri, row_data in enumerate(rows):
                y = margin + 44 + ri * row_h
                if y + row_h > A4_H - margin:
                    page = doc.new_page()
                    y = margin
                x = margin
                for ci, val in enumerate(row_data):
                    cell_x = x + ci * max_text_w
                    # 表头加背景色
                    if ri == 0:
                        page.draw_rect(fitz.Rect(cell_x - 1, y - 1, cell_x + max_text_w - 1, y + row_h - 1),
                                      color=(0.5, 0.5, 0.5), width=0.3, fill=(0.85, 0.85, 0.85))
                    else:
                        page.draw_rect(fitz.Rect(cell_x - 1, y - 1, cell_x + max_text_w - 1, y + row_h - 1),
                                      color=(0.7, 0.7, 0.7), width=0.2)
                    text = str(val) if val is not None else ""
                    page.insert_text((cell_x + 2, y + 12), text[:30], fontsize=8, fontname="hebo", fontfile=None)

        out_path = _output_path(fid)
        doc.save(out_path, garbage=4, deflate=True)
        doc.close()
        return send_file(out_path, as_attachment=True,
                        download_name=f"{Path(f.filename).stem}.pdf")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/convert/to-pdf/ppt", methods=["POST"])
def ppt_to_pdf():
    """PPT (.pptx) 转 PDF —— 每张幻灯片一页"""
    f = request.files.get("file")
    if not f: return json_err("请上传PPT文件")
    fpath, fid = save_upload(f)
    try:
        from pptx import Presentation
        prs = Presentation(fpath)
        doc = fitz.open()
        A4_W, A4_H = 595, 842
        margin = 48

        for si, slide in enumerate(prs.slides):
            page = doc.new_page(width=A4_W, height=A4_H)
            y = margin
            page.insert_text((margin, y - 10), f"Slide {si + 1}", fontsize=12, color=(0.4, 0.4, 0.4), fontname="hebo", fontfile=None)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text: continue
                        y += 16
                        if y > A4_H - margin:
                            page = doc.new_page(width=A4_W, height=A4_H)
                            y = margin
                        fs = 10
                        try:
                            run_fs = para.runs[0].font.size
                            if run_fs: fs = run_fs / 12700  # EMU to pt
                        except: pass
                        page.insert_text((margin, y), text[:120], fontsize=min(fs, 18), fontname="hebo", fontfile=None)
                elif shape.shape_type == 13:  # 图片
                    try:
                        img_data = shape.image.blob
                        pix = fitz.open("png", img_data)
                        y += 20
                        if y + 140 > A4_H - margin:
                            page = doc.new_page(width=A4_W, height=A4_H)
                            y = margin
                        page.insert_image(fitz.Rect(margin + 20, y, margin + 200, y + 140), stream=img_data)
                        y += 150
                    except: pass

        out_path = _output_path(fid)
        doc.save(out_path, garbage=4, deflate=True)
        doc.close()
        return send_file(out_path, as_attachment=True,
                        download_name=f"{Path(f.filename).stem}.pdf")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/convert/to-pdf/html", methods=["POST"])
def html_to_pdf():
    """HTML 转 PDF —— 提取文本内容渲染到A4页面"""
    f = request.files.get("file")
    if not f: return json_err("请上传HTML文件")
    fpath, fid = save_upload(f)
    try:
        with open(fpath, "r", encoding="utf-8", errors="ignore") as fh:
            html_content = fh.read()

        # 尝试 weasyprint（需要 GTK 运行时，Windows 上通常不可用）
        try:
            from weasyprint import HTML as WHTML
            out_path = _output_path(fid)
            WHTML(string=html_content).write_pdf(out_path)
            return send_file(out_path, as_attachment=True,
                            download_name=f"{Path(f.filename).stem}.pdf")
        except Exception:
            pass  # weasyprint 不可用，使用以下 fallback

        # Fallback: HTML 文本提取 + fitz 渲染
        from html.parser import HTMLParser

        class HtmlToPdfParser(HTMLParser):
            def __init__(self):
                super().__init__()
                self.blocks = []  # (tag_type, text)
                self.skip = False
                self._cur = ""
                self._cur_tag = "p"

            def handle_starttag(self, tag, attrs):
                tag = tag.lower()
                if tag in ("script", "style", "head"):
                    self.skip = True
                elif tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                    self._flush()
                    self._cur_tag = tag
                elif tag in ("p", "div", "br", "li", "td", "th", "tr", "pre"):
                    self._flush()
                    self._cur_tag = tag
                elif tag == "hr":
                    self._flush()
                    self.blocks.append(("hr", ""))

            def handle_endtag(self, tag):
                tag = tag.lower()
                if tag in ("script", "style", "head"):
                    self.skip = False
                elif tag in ("h1", "h2", "h3", "h4", "h5", "h6", "p", "div", "li", "td", "th"):
                    self._flush()
                    self._cur_tag = "p"

            def _flush(self):
                t = self._cur.strip()
                if t:
                    self.blocks.append((self._cur_tag, t))
                self._cur = ""

            def handle_data(self, data):
                if not self.skip:
                    self._cur += data

        parser = HtmlToPdfParser()
        parser.feed(html_content)
        parser._flush()

        doc = fitz.open()
        A4_W, A4_H = 595, 842
        margin = 50
        y = margin
        page = doc.new_page()
        h_sizes = {"h1": 20, "h2": 17, "h3": 14, "h4": 12, "h5": 11, "h6": 10}

        for tag, text in parser.blocks:
            if tag == "hr":
                y += 6
                if y > A4_H - margin:
                    page = doc.new_page()
                    y = margin
                page.draw_line((margin, y), (A4_W - margin, y), color=(0.7, 0.7, 0.7), width=0.5)
                y += 10
                continue

            fs = h_sizes.get(tag, 10)
            lh = fs + 8
            if y + lh > A4_H - margin:
                page = doc.new_page()
                y = margin
            indent = 20 if tag in ("li",) else 0
            if tag in ("pre",):
                for pline in text.split("\n"):
                    if y + 14 > A4_H - margin:
                        page = doc.new_page()
                        y = margin
                    page.insert_text((margin + indent + 10, y), pline[:100], fontsize=9,
                                    fontname="hebo", fontfile=None)
                    y += 14
                y += 4
            else:
                page.insert_text((margin + indent, y), text[:120], fontsize=fs,
                                fontname="hebo", fontfile=None)
                y += lh

        out_path = _output_path(fid)
        doc.save(out_path, garbage=4, deflate=True)
        doc.close()
        return send_file(out_path, as_attachment=True,
                        download_name=f"{Path(f.filename).stem}.pdf")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/convert/to-pdf/markdown", methods=["POST"])
def markdown_to_pdf():
    """Markdown (.md) 转 PDF —— 解析标题/列表/表格，分级字体渲染"""
    f = request.files.get("file")
    if not f: return json_err("请上传Markdown文件")
    fpath, fid = save_upload(f)
    try:
        with open(fpath, "r", encoding="utf-8", errors="ignore") as fh:
            md_text = fh.read()

        doc = fitz.open()
        A4_W, A4_H = 595, 842
        margin = 50
        y = margin
        page = doc.new_page()

        in_code_block = False

        for line in md_text.split("\n"):
            # 代码块切换
            if line.startswith("```"):
                in_code_block = not in_code_block
                continue
            if in_code_block:
                if y + 14 > A4_H - margin:
                    page = doc.new_page()
                    y = margin
                page.draw_rect(fitz.Rect(margin - 2, y - 10, A4_W - margin + 2, y + 4),
                               fill=(0.94, 0.94, 0.94), color=None)
                page.insert_text((margin + 4, y), line[:100], fontsize=9,
                                fontname="hebo", fontfile=None, color=(0.2, 0.2, 0.2))
                y += 14
                continue

            stripped = line.strip()
            if not stripped:
                y += 8
                continue

            # 标题级别
            if stripped.startswith("#"):
                level = len(stripped) - len(stripped.lstrip("#"))
                level = min(level, 4)
                text = stripped.lstrip("#").strip()
                fs_map = {1: 22, 2: 17, 3: 14, 4: 12}
                fs = fs_map.get(level, 12)
                y += 4
                if y + fs + 10 > A4_H - margin:
                    page = doc.new_page()
                    y = margin
                # 标题下划线
                if level <= 2:
                    page.draw_line((margin, y + fs + 2), (A4_W - margin, y + fs + 2),
                                   color=(0.8, 0.8, 0.8), width=0.5)
                page.insert_text((margin, y), text[:100], fontsize=fs,
                                fontname="hebo", fontfile=None)
                y += fs + 12
                continue

            # 列表项
            if stripped.startswith(("- ", "* ", "+ ")):
                text = "  • " + stripped[2:].strip()
                if y + 16 > A4_H - margin:
                    page = doc.new_page()
                    y = margin
                page.insert_text((margin + 8, y), text[:120], fontsize=10,
                                fontname="hebo", fontfile=None)
                y += 16
                continue

            # 有序列表
            import re
            m = re.match(r"^(\d+)\. (.+)$", stripped)
            if m:
                text = f"  {m.group(1)}. {m.group(2)}"
                if y + 16 > A4_H - margin:
                    page = doc.new_page()
                    y = margin
                page.insert_text((margin + 8, y), text[:120], fontsize=10,
                                fontname="hebo", fontfile=None)
                y += 16
                continue

            # 水平线
            if re.match(r"^[-*=_]{3,}$", stripped):
                if y + 10 > A4_H - margin:
                    page = doc.new_page()
                    y = margin
                page.draw_line((margin, y), (A4_W - margin, y), color=(0.7, 0.7, 0.7), width=0.5)
                y += 14
                continue

            # 表格行
            if stripped.startswith("|") and stripped.endswith("|"):
                cells = [c.strip() for c in stripped.strip("|").split("|")]
                if all(re.match(r"^[-:]+$", c) for c in cells if c):
                    continue  # 表头分隔行
                ncols = max(len(cells), 1)
                col_w = (A4_W - 2 * margin) / ncols
                if y + 20 > A4_H - margin:
                    page = doc.new_page()
                    y = margin
                for i, cell in enumerate(cells):
                    cx = margin + i * col_w
                    page.draw_rect(fitz.Rect(cx, y - 2, cx + col_w - 1, y + 16),
                                   color=(0.7, 0.7, 0.7), width=0.3)
                    page.insert_text((cx + 3, y + 10), cell[:30], fontsize=9,
                                    fontname="hebo", fontfile=None)
                y += 20
                continue

            # 普通段落（去除 Markdown 内联标记显示）
            text = re.sub(r"\*\*(.+?)\*\*", r"\1", stripped)
            text = re.sub(r"\*(.+?)\*", r"\1", text)
            text = re.sub(r"`(.+?)`", r"\1", text)
            text = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", text)

            if y + 18 > A4_H - margin:
                page = doc.new_page()
                y = margin
            page.insert_text((margin, y), text[:140], fontsize=10,
                            fontname="hebo", fontfile=None)
            y += 18

        out_path = _output_path(fid)
        doc.save(out_path, garbage=4, deflate=True)
        doc.close()
        return send_file(out_path, as_attachment=True,
                        download_name=f"{Path(f.filename).stem}.pdf")
    except Exception as e:
        return json_err(str(e))
        return json_err(str(e))


@app.route("/api/pdf/convert/to-pdf/image", methods=["POST"])
def image_to_pdf():
    """图片转 PDF —— 支持多张图片，自动分页并保持比例"""
    files = request.files.getlist("images")
    if not files:
        img = request.files.get("file") or request.files.get("image")
        if not img: return json_err("请上传图片")
        files = [img]
    try:
        A4_W, A4_H = 595, 842
        margin = 20
        doc = fitz.open()
        fid = None

        for f in files:
            fpath, _fid = save_upload(f)
            fid = _fid
            pil_img = Image.open(fpath)
            if pil_img.mode != "RGB":
                pil_img = pil_img.convert("RGB")
            pw, ph = pil_img.size
            max_w = A4_W - 2 * margin
            max_h = A4_H - 2 * margin
            scale = min(max_w / pw, max_h / ph)
            iw, ih = pw * scale, ph * scale
            cx, cy = (A4_W - iw) / 2, (A4_H - ih) / 2

            page = doc.new_page(width=A4_W, height=A4_H)
            img_buf = io.BytesIO()
            pil_img.save(img_buf, "JPEG", quality=92)
            page.insert_image(fitz.Rect(cx, cy, cx + iw, cy + ih), stream=img_buf.getvalue())
            pil_img.close()

        out_path = _output_path(fid)
        doc.save(out_path, garbage=4, deflate=True)
        doc.close()
        return send_file(out_path, as_attachment=True,
                        download_name="images_to_pdf.pdf")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/convert/to-pdf/txt", methods=["POST"])
def txt_to_pdf():
    """文本 (.txt) 转 PDF —— 智能分页，支持中文"""
    f = request.files.get("file")
    if not f: return json_err("请上传文本文件")
    fpath, fid = save_upload(f)
    try:
        with open(fpath, "r", encoding="utf-8", errors="replace") as fh:
            content = fh.read()

        doc = fitz.open()
        A4_W, A4_H = 595, 842
        margin = 50
        line_height = 18
        chars_per_line = 72  # fontsize=11 大约容纳
        lines_per_page = int((A4_H - 2 * margin) / line_height)

        lines = []
        for paragraph in content.split("\n"):
            if not paragraph.strip():
                lines.append("")
                continue
            # 长行自动换行
            while len(paragraph) > chars_per_line:
                lines.append(paragraph[:chars_per_line])
                paragraph = paragraph[chars_per_line:]
            lines.append(paragraph)

        page = doc.new_page()
        y = margin
        for i, line in enumerate(lines):
            if i > 0 and i % lines_per_page == 0:
                page = doc.new_page()
                y = margin
            page.insert_text((margin, y), line[:chars_per_line], fontsize=11, fontname="hebo", fontfile=None)
            y += line_height

        out_path = _output_path(fid)
        doc.save(out_path, garbage=4, deflate=True)
        doc.close()
        return send_file(out_path, as_attachment=True,
                        download_name=f"{Path(f.filename).stem}.pdf")
    except Exception as e:
        return json_err(str(e))


# ---- 旧接口保持向后兼容 ----
@app.route("/api/pdf/convert/office-to-pdf", methods=["POST"])
def office_to_pdf():
    """(兼容旧版) Office/图片/文本转PDF — 自动识别格式"""
    f = request.files.get("file")
    if not f: return json_err("请上传文件")
    fpath, fid = save_upload(f)
    ext = Path(f.filename).suffix.lower()
    try:
        if ext in [".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff", ".webp"]:
            # 图片转PDF
            img = Image.open(fpath)
            if img.mode != "RGB":
                img = img.convert("RGB")
            pdf_bytes = io.BytesIO()
            img.save(pdf_bytes, "PDF")
            pdf_bytes.seek(0)
            doc = fitz.open("pdf", pdf_bytes.read())
            out_path = _output_path(fid)
            doc.save(out_path)
            doc.close()
            return send_file(out_path, as_attachment=True,
                           download_name=f"{Path(f.filename).stem}.pdf")
        elif ext == ".txt":
            # 委托到专用的 txt→pdf 端点
            return txt_to_pdf()
        elif ext in [".docx", ".doc"]:
            return word_to_pdf()
        elif ext in [".xlsx", ".xls"]:
            return excel_to_pdf()
        elif ext in [".pptx", ".ppt"]:
            return ppt_to_pdf()
        elif ext in [".html", ".htm"]:
            return html_to_pdf()
        elif ext == ".md":
            return markdown_to_pdf()
        else:
            return json_err(f"暂不支持 {ext} 格式转PDF，请使用专用转换入口")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/convert/to-txt", methods=["POST"])
def pdf_to_txt():
    """PDF转TXT"""
    f = request.files.get("file")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        all_text = []
        for i in range(doc.page_count):
            all_text.append(doc[i].get_text("text"))
        doc.close()
        out_path = _output_path(fid, ".txt")
        with open(out_path, "w", encoding="utf-8") as fh:
            fh.write("\n".join(all_text))
        return send_file(out_path, as_attachment=True, download_name=f"{Path(f.filename).stem}.txt")
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 5. PDF 压缩
# ============================================================

@app.route("/api/pdf/compress", methods=["POST"])
def pdf_compress():
    """压缩PDF（轻微/中等/大幅），返回压缩文件+统计"""
    f = request.files.get("file")
    level = request.form.get("level", "medium")  # light, medium, heavy
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        orig_size = os.path.getsize(fpath)
        doc = fitz.open(fpath)

        # 压缩级别映射（JPEG 质量）
        quality_map = {"light": 85, "medium": 50, "heavy": 20}
        jpeg_quality = quality_map.get(level, 50)
        dpi_map = {"light": 200, "medium": 150, "heavy": 100}
        render_dpi = dpi_map.get(level, 150)

        doc_new = fitz.open()
        for page in doc:
            pix = page.get_pixmap(dpi=render_dpi)
            img_data = pix.tobytes("jpeg", jpeg_quality)
            new_page = doc_new.new_page(width=page.rect.width, height=page.rect.height)
            new_page.insert_image(new_page.rect, stream=img_data)

        out_path = _output_path(fid)
        doc_new.save(out_path, garbage=4, deflate=True)
        doc.close()
        doc_new.close()

        new_size = os.path.getsize(out_path)
        # 如果压缩后更大，退回原始文件（文本型PDF不需要图像重压缩）
        if new_size >= orig_size:
            shutil.copy2(fpath, out_path)
            new_size = orig_size

        return send_file(out_path, as_attachment=True,
                        download_name=f"compressed_{f.filename}")
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 6. PDF 合并/拆分/提取
# ============================================================

@app.route("/api/pdf/merge", methods=["POST"])
def pdf_merge():
    """合并多个PDF"""
    files = request.files.getlist("files")
    if len(files) < 2:
        return json_err("请至少上传2个PDF文件")
    try:
        doc_out = fitz.open()
        for f in files:
            fpath, _ = save_upload(f)
            doc_src = fitz.open(fpath)
            doc_out.insert_pdf(doc_src)
            doc_src.close()
        out_path = _output_path(gen_id())
        doc_out.save(out_path)
        doc_out.close()
        return send_file(out_path, as_attachment=True, download_name="merged.pdf")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/split", methods=["POST"])
def pdf_split():
    """拆分PDF"""
    f = request.files.get("file")
    mode = request.form.get("mode", "all")  # all, range, every
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        total = doc.page_count

        if mode == "every":
            # 每N页拆分
            n = int(request.form.get("n", 1))
            zip_buf = io.BytesIO()
            with zipfile.ZipFile(zip_buf, "w") as zf:
                for start in range(0, total, n):
                    end = min(start + n, total)
                    sub_doc = fitz.open()
                    sub_doc.insert_pdf(doc, from_page=start, to_page=end - 1)
                    buf = io.BytesIO()
                    sub_doc.save(buf)
                    sub_doc.close()
                    zf.writestr(f"split_{start+1}-{end}.pdf", buf.getvalue())
            zip_buf.seek(0)
            doc.close()
            return send_file(zip_buf, mimetype="application/zip",
                           download_name=f"{Path(f.filename).stem}_split.zip")
        elif mode == "range":
            # 按页码范围拆分，支持 "1-3,5-7" 和 JSON [[1,3],[5,7]] 两种格式
            ranges_raw = request.form.get("ranges", "[]")
            try:
                ranges = json.loads(ranges_raw)
            except (json.JSONDecodeError, ValueError):
                # 解析 "1-3,5-7" 格式
                ranges = []
                for part in ranges_raw.split(","):
                    part = part.strip()
                    if "-" in part:
                        a, b = part.split("-", 1)
                        ranges.append([int(a.strip()), int(b.strip())])
                    elif part.isdigit():
                        n = int(part)
                        ranges.append([n, n])
            zip_buf = io.BytesIO()
            with zipfile.ZipFile(zip_buf, "w") as zf:
                for rg in ranges:
                    start, end = rg[0] - 1, rg[1] - 1
                    sub_doc = fitz.open()
                    sub_doc.insert_pdf(doc, from_page=start, to_page=end)
                    buf = io.BytesIO()
                    sub_doc.save(buf)
                    sub_doc.close()
                    zf.writestr(f"pages_{rg[0]}-{rg[1]}.pdf", buf.getvalue())
            zip_buf.seek(0)
            doc.close()
            return send_file(zip_buf, mimetype="application/zip",
                           download_name=f"{Path(f.filename).stem}_split_range.zip")
        else:
            # 每页单独拆
            zip_buf = io.BytesIO()
            with zipfile.ZipFile(zip_buf, "w") as zf:
                for i in range(total):
                    sub_doc = fitz.open()
                    sub_doc.insert_pdf(doc, from_page=i, to_page=i)
                    buf = io.BytesIO()
                    sub_doc.save(buf)
                    sub_doc.close()
                    zf.writestr(f"page_{i+1}.pdf", buf.getvalue())
            zip_buf.seek(0)
            doc.close()
            return send_file(zip_buf, mimetype="application/zip",
                           download_name=f"{Path(f.filename).stem}_split_all.zip")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/extract", methods=["POST"])
def pdf_extract_pages():
    """提取指定页面"""
    f = request.files.get("file")
    pages_str = request.form.get("pages", "")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        pages = [int(p.strip()) - 1 for p in pages_str.split(",") if p.strip()]
        doc = fitz.open(fpath)
        doc_out = fitz.open()
        for p in pages:
            if 0 <= p < doc.page_count:
                doc_out.insert_pdf(doc, from_page=p, to_page=p)
        out_path = _output_path(fid)
        doc_out.save(out_path)
        doc.close()
        doc_out.close()
        return send_file(out_path, as_attachment=True, download_name=f"extracted_{f.filename}")
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 7. OCR 识别
# ============================================================

@app.route("/api/pdf/ocr", methods=["POST"])
def pdf_ocr():
    """OCR识别"""
    f = request.files.get("file")
    lang = request.form.get("lang", "chi_sim+eng")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        import pytesseract
        doc = fitz.open(fpath)
        results = []
        for i in range(doc.page_count):
            page = doc[i]
            pix = page.get_pixmap(dpi=300)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            text = pytesseract.image_to_string(img, lang=lang)
            results.append({"page": i + 1, "text": text})
        doc.close()
        return json_ok(results)
    except ImportError:
        return json_err("OCR引擎未安装，请安装Tesseract-OCR")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/ocr/make-searchable", methods=["POST"])
def pdf_ocr_searchable():
    """将扫描PDF转为可搜索PDF"""
    f = request.files.get("file")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        import pytesseract
        doc = fitz.open(fpath)
        doc_out = fitz.open()
        for i in range(doc.page_count):
            page = doc[i]
            pix = page.get_pixmap(dpi=200)
            img_data = pix.tobytes("png")
            new_page = doc_out.new_page(width=page.rect.width, height=page.rect.height)
            new_page.insert_image(new_page.rect, stream=img_data)
            # OCR文字层
            img = Image.open(io.BytesIO(img_data))
            ocr_data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT, lang="chi_sim+eng")
            for j in range(len(ocr_data["text"])):
                if ocr_data["text"][j].strip():
                    x, y, w, h = ocr_data["left"][j], ocr_data["top"][j], ocr_data["width"][j], ocr_data["height"][j]
                    scale = page.rect.width / img.width
                    new_page.insert_text((x * scale, y * scale + h * scale * 0.8),
                                         ocr_data["text"][j], fontsize=1, color=(1, 1, 1))
        out_path = _output_path(fid)
        doc_out.save(out_path)
        doc.close()
        doc_out.close()
        return send_file(out_path, as_attachment=True, download_name=f"ocr_{f.filename}")
    except ImportError:
        return json_err("OCR引擎未安装")
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 8. 水印
# ============================================================

@app.route("/api/pdf/watermark", methods=["POST"])
def pdf_add_watermark():
    """添加水印"""
    f = request.files.get("file")
    wm_type = request.form.get("type", "text")  # text or image
    wm_text = request.form.get("text", "保密")
    opacity = float(request.form.get("opacity", 0.3))
    rotation = float(request.form.get("rotation", 45))
    font_size = int(request.form.get("font_size", 60))

    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        for page in doc:
            rect = page.rect
            c = rect.width / 2
            d = rect.height / 2
            if wm_type == "text":
                # 平铺水印
                for x in range(-100, int(rect.width) + 200, 300):
                    for y in range(-100, int(rect.height) + 200, 200):
                        page.insert_text(
                            (x, y), wm_text,
                            fontname="china-s",
                            fontsize=font_size,
                            color=(0.5, 0.5, 0.5),
                            alpha=opacity,
                            rotate=rotation,
                        )
            elif wm_type == "image":
                wm_file = request.files.get("watermark_image")
                if wm_file:
                    wm_img = Image.open(wm_file)
                    wm_bytes = io.BytesIO()
                    wm_img.save(wm_bytes, format="PNG")
                    wm_bytes.seek(0)
                    # 居中放置
                    w, h = wm_img.width * 0.5, wm_img.height * 0.5
                    page.insert_image(
                        fitz.Rect(c - w/2, d - h/2, c + w/2, d + h/2),
                        stream=wm_bytes.getvalue(),
                        alpha=int(opacity * 255),
                    )
        out_path = _output_path(fid)
        doc.save(out_path)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name=f"watermark_{f.filename}")
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 9. 电子签名
# ============================================================

@app.route("/api/pdf/signature", methods=["POST"])
def pdf_add_signature():
    """添加电子签名"""
    f = request.files.get("file")
    sig_type = request.form.get("type", "draw")  # draw, text, image
    page_num = int(request.form.get("page", 1))
    x = float(request.form.get("x", 100))
    y = float(request.form.get("y", 100))
    w = float(request.form.get("w", 150))
    h = float(request.form.get("h", 60))

    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        page = doc[page_num - 1]
        rect = fitz.Rect(x, y, x + w, y + h)

        if sig_type == "text":
            text = request.form.get("text", "")
            font_size = int(request.form.get("font_size", 24))
            page.insert_text((x, y + h * 0.7), text, fontname="china-s",
                           fontsize=font_size, color=(0, 0, 0.7))
        elif sig_type == "image":
            sig_file = request.files.get("signature_image")
            if sig_file:
                page.insert_image(rect, stream=sig_file.read())
        elif sig_type == "draw":
            # 手写签名：渲染前端传来的SVG/base64
            sig_data = request.form.get("sig_data", "")
            if sig_data.startswith("data:image"):
                import re
                img_b64 = re.sub(r"^data:image/\w+;base64,", "", sig_data)
                img_bytes = base64.b64decode(img_b64)
                page.insert_image(rect, stream=img_bytes)

        out_path = _output_path(fid)
        doc.save(out_path)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name=f"signed_{f.filename}")
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 10. 加密/安全
# ============================================================

@app.route("/api/pdf/encrypt", methods=["POST"])
def pdf_encrypt():
    """加密PDF"""
    f = request.files.get("file")
    password = request.form.get("password", "")
    if not f or not password:
        return json_err("请上传文件并设置密码")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        out_path = _output_path(fid)
        doc.save(out_path, encryption=fitz.PDF_ENCRYPT_AES_256, owner_pw=password, user_pw=password)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name=f"encrypted_{f.filename}")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/decrypt", methods=["POST"])
def pdf_decrypt():
    """解密PDF"""
    f = request.files.get("file")
    password = request.form.get("password", "")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        if doc.needs_pass:
            ok = doc.authenticate(password)
            if not ok:
                return json_err("密码错误")
        out_path = _output_path(fid)
        doc.save(out_path)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name=f"decrypted_{f.filename}")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/protect", methods=["POST"])
def pdf_protect_permissions():
    """设置PDF权限"""
    f = request.files.get("file")
    password = request.form.get("password", "")
    permissions = json.loads(request.form.get("permissions", "{}"))
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        perm = (
            fitz.PDF_PERM_PRINT * permissions.get("print", 0) |
            fitz.PDF_PERM_MODIFY * permissions.get("modify", 0) |
            fitz.PDF_PERM_COPY * permissions.get("copy", 0) |
            fitz.PDF_PERM_ANNOTATE * permissions.get("annotate", 0)
        )
        out_path = _output_path(fid)
        doc.save(out_path, encryption=fitz.PDF_ENCRYPT_AES_256, owner_pw=password, permissions=perm)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name=f"protected_{f.filename}")
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 11. 页面管理
# ============================================================

@app.route("/api/pdf/pages/reorder", methods=["POST"])
def pdf_reorder():
    """页面重新排序"""
    f = request.files.get("file")
    order_str = request.form.get("order", "[]")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    order = json.loads(order_str)  # [3,1,2,4,...]
    try:
        doc = fitz.open(fpath)
        doc_out = fitz.open()
        for p in order:
            if 0 <= p - 1 < doc.page_count:
                doc_out.insert_pdf(doc, from_page=p - 1, to_page=p - 1)
        out_path = _output_path(fid)
        doc_out.save(out_path)
        doc.close()
        doc_out.close()
        return send_file(out_path, as_attachment=True, download_name=f"reordered_{f.filename}")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/pages/rotate", methods=["POST"])
def pdf_rotate():
    """旋转页面"""
    f = request.files.get("file")
    pages_str = request.form.get("pages", "")
    angle = int(request.form.get("angle", 90))
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        pages = [int(p.strip()) - 1 for p in pages_str.split(",") if p.strip()]
        if not pages:
            pages = list(range(doc.page_count))
        for p in pages:
            if 0 <= p < doc.page_count:
                doc[p].set_rotation(doc[p].rotation + angle)
        out_path = _output_path(fid)
        doc.save(out_path)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name=f"rotated_{f.filename}")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/pages/crop", methods=["POST"])
def pdf_crop():
    """裁剪页面"""
    f = request.files.get("file")
    crop_data = json.loads(request.form.get("crop", "{}"))
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        for page_num_str, rect_data in crop_data.items():
            p = int(page_num_str) - 1
            if 0 <= p < doc.page_count:
                page = doc[p]
                page.set_cropbox(fitz.Rect(rect_data[0], rect_data[1], rect_data[2], rect_data[3]))
        out_path = _output_path(fid)
        doc.save(out_path)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name=f"cropped_{f.filename}")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/pages/delete", methods=["POST"])
def pdf_delete_pages():
    """删除页面"""
    f = request.files.get("file")
    pages_str = request.form.get("pages", "")
    if not f:
        return json_err("请上传文件")
    fpath, fid = save_upload(f)
    try:
        delete_set = set(int(p.strip()) - 1 for p in pages_str.split(",") if p.strip())
        doc = fitz.open(fpath)
        doc_out = fitz.open()
        for i in range(doc.page_count):
            if i not in delete_set:
                doc_out.insert_pdf(doc, from_page=i, to_page=i)
        out_path = _output_path(fid)
        doc_out.save(out_path)
        doc.close()
        doc_out.close()
        return send_file(out_path, as_attachment=True, download_name=f"deleted_{f.filename}")
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 12. PDF 创建
# ============================================================

@app.route("/api/pdf/create", methods=["POST"])
def pdf_create():
    """创建空白PDF"""
    width = float(request.form.get("width", 595))  # A4
    height = float(request.form.get("height", 842))
    pages = int(request.form.get("pages", 1))
    try:
        doc = fitz.open()
        for _ in range(pages):
            doc.new_page(width=width, height=height)
        out_path = _output_path(gen_id())
        doc.save(out_path)
        doc.close()
        return send_file(out_path, as_attachment=True, download_name="new_document.pdf")
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 13. 搜索
# ============================================================

@app.route("/api/pdf/search", methods=["POST"])
def pdf_search():
    """搜索PDF文本"""
    f = request.files.get("file")
    keyword = request.form.get("keyword", "")
    if not f or not keyword:
        return json_err("请上传文件并输入关键词")
    fpath, fid = save_upload(f)
    try:
        doc = fitz.open(fpath)
        results = []
        for i in range(doc.page_count):
            page = doc[i]
            areas = page.search_for(keyword)
            for area in areas:
                results.append({
                    "page": i + 1,
                    "x": area.x0, "y": area.y0,
                    "w": area.width, "h": area.height,
                })
        doc.close()
        return json_ok({"keyword": keyword, "count": len(results), "results": results})
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 14. 打印支持
# ============================================================

@app.route("/api/pdf/print/layout", methods=["POST"])
def pdf_print_layout():
    """生成打印排版（多页排版到一页，保持原图比例）支持多文件"""
    pages_per_sheet = int(request.form.get("per_sheet", 2))

    # 支持多文件上传（files[] 或多个 file 参数）
    files = request.files.getlist("files")
    if not files:
        f = request.files.get("file")
        if not f:
            return json_err("请上传文件")
        files = [f]

    # 先合并所有文件为一个文档
    try:
        doc = fitz.open()
        fid = None
        for f in files:
            fpath, _fid = save_upload(f)
            fid = _fid
            src = fitz.open(fpath)
            doc.insert_pdf(src)
            src.close()
    except Exception as e:
        return json_err(str(e))

    try:
        doc_out = fitz.open()

        # A4 尺寸
        A4_W, A4_H = 595, 842

        if pages_per_sheet == 2:
            # 上下布局，保持比例
            gap = 6  # 间距
            line_y = A4_H / 2
            top_h = line_y - gap / 2
            bot_h = A4_H - line_y - gap / 2

            for i in range(0, doc.page_count, 2):
                page_a4 = doc_out.new_page(width=A4_W, height=A4_H)

                # 上半部分
                _draw_page_scaled(page_a4, doc[i],
                    fitz.Rect(0, 0, A4_W, top_h), margin=8)

                # 裁剪虚线
                _draw_dash_line(page_a4, (0, line_y), (A4_W, line_y))

                # 下半部分
                if i + 1 < doc.page_count:
                    _draw_page_scaled(page_a4, doc[i + 1],
                        fitz.Rect(0, line_y + gap / 2, A4_W, A4_H), margin=8)

            out_path = _output_path(fid)
            doc_out.save(out_path)
            doc_out.close()
            doc.close()
            return send_file(out_path, as_attachment=True,
                           download_name=f"print_2up.pdf")

        elif pages_per_sheet == 4:
            # 2x2 网格布局
            half_w = A4_W / 2
            half_h = A4_H / 2
            margin = 4
            positions = [
                fitz.Rect(margin, margin, half_w - margin, half_h - margin),
                fitz.Rect(half_w + margin, margin, A4_W - margin, half_h - margin),
                fitz.Rect(margin, half_h + margin, half_w - margin, A4_H - margin),
                fitz.Rect(half_w + margin, half_h + margin, A4_W - margin, A4_H - margin),
            ]

            for i in range(0, doc.page_count, 4):
                page_a4 = doc_out.new_page(width=A4_W, height=A4_H)
                for j, rect in enumerate(positions):
                    if i + j < doc.page_count:
                        _draw_page_scaled(page_a4, doc[i + j], rect, margin=0)

            out_path = _output_path(fid)
            doc_out.save(out_path)
            doc_out.close()
            doc.close()
            return send_file(out_path, as_attachment=True,
                           download_name=f"print_4up.pdf")

        elif pages_per_sheet == 6:
            # 2x3 网格（适合PPT打印）
            third_h = A4_H / 3
            half_w = A4_W / 2
            margin = 3
            positions = []
            for row in range(3):
                for col in range(2):
                    positions.append(fitz.Rect(
                        col * half_w + margin, row * third_h + margin,
                        (col + 1) * half_w - margin, (row + 1) * third_h - margin
                    ))

            for i in range(0, doc.page_count, 6):
                page_a4 = doc_out.new_page(width=A4_W, height=A4_H)
                for j, rect in enumerate(positions):
                    if i + j < doc.page_count:
                        _draw_page_scaled(page_a4, doc[i + j], rect, margin=0)

            out_path = _output_path(fid)
            doc_out.save(out_path)
            doc_out.close()
            doc.close()
            return send_file(out_path, as_attachment=True,
                           download_name=f"print_6up.pdf")

        elif pages_per_sheet == 9:
            # 3x3 网格
            third_w = A4_W / 3
            third_h = A4_H / 3
            margin = 2
            positions = []
            for row in range(3):
                for col in range(3):
                    positions.append(fitz.Rect(
                        col * third_w + margin, row * third_h + margin,
                        (col + 1) * third_w - margin, (row + 1) * third_h - margin
                    ))

            for i in range(0, doc.page_count, 9):
                page_a4 = doc_out.new_page(width=A4_W, height=A4_H)
                for j, rect in enumerate(positions):
                    if i + j < doc.page_count:
                        _draw_page_scaled(page_a4, doc[i + j], rect, margin=0)

            out_path = _output_path(fid)
            doc_out.save(out_path)
            doc_out.close()
            doc.close()
            return send_file(out_path, as_attachment=True,
                           download_name=f"print_9up.pdf")

        doc.close()
        doc_out.close()
        return json_err("支持 2/4/6/9 合1排版")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/print/direct", methods=["POST"])
def pdf_print_direct():
    """直接打印：合并多个PDF为一个，供浏览器打印对话框使用"""
    files = request.files.getlist("files")
    if not files:
        f = request.files.get("file")
        if not f:
            return json_err("请上传文件")
        files = [f]

    try:
        doc = fitz.open()
        fid = None
        for f in files:
            fpath, _fid = save_upload(f)
            fid = _fid
            src = fitz.open(fpath)
            doc.insert_pdf(src)
            src.close()

        out_path = _output_path(fid)
        doc.save(out_path, garbage=4, deflate=True)
        doc.close()
        return send_file(out_path, as_attachment=True,
                       download_name="print_direct.pdf")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/pdf/print/images", methods=["POST"])
def pdf_print_images():
    """图片合并与多合一排版：接收多张图片，按指定布局生成打印PDF"""
    pages_per_sheet = int(request.form.get("per_sheet", 1))
    images = request.files.getlist("images")
    if not images:
        # 兼容单文件字段
        img = request.files.get("image") or request.files.get("file")
        if not img:
            return json_err("请上传图片文件")
        images = [img]

    try:
        from PIL import Image as PILImage
    except ImportError:
        return json_err("服务端Pillow库未安装")

    A4_W, A4_H = 595, 842
    doc = fitz.open()
    total_images = len(images)
    fid = None

    # 保存所有图片到临时目录
    img_paths = []
    for img_file in images:
        fpath, _fid = save_upload(img_file)
        fid = _fid
        img_paths.append(fpath)

    if pages_per_sheet == 1:
        # 每页一张图，自适应A4
        for fpath in img_paths:
            pim = PILImage.open(fpath).convert("RGB")
            pw, ph = pim.size
            scale = min(A4_W / pw, A4_H / ph)
            iw, ih = pw * scale, ph * scale
            cx, cy = (A4_W - iw) / 2, (A4_H - ih) / 2

            page = doc.new_page(width=A4_W, height=A4_H)
            buf = io.BytesIO()
            pim.save(buf, "JPEG", quality=90)
            page.insert_image(fitz.Rect(cx, cy, cx + iw, cy + ih), stream=buf.getvalue())
            pim.close()

    elif pages_per_sheet in (2, 4, 6, 9):
        # 定义网格
        if pages_per_sheet == 2:
            cols, rows = 1, 2
        elif pages_per_sheet == 4:
            cols, rows = 2, 2
        elif pages_per_sheet == 6:
            cols, rows = 2, 3
        else:  # 9
            cols, rows = 3, 3

        margin = 8
        cell_w = (A4_W - margin * 2) / cols
        cell_h = (A4_H - margin * 2) / rows
        slots_per_page = cols * rows
        img_idx = 0

        while img_idx < total_images:
            page = doc.new_page(width=A4_W, height=A4_H)
            for slot in range(slots_per_page):
                if img_idx >= total_images:
                    break
                row, col = divmod(slot, cols)
                cell_x = margin + col * cell_w
                cell_y = margin + row * cell_h

                pim = PILImage.open(img_paths[img_idx]).convert("RGB")
                pw, ph = pim.size
                inner_m = 4  # 图片之间的间距
                inner_w = cell_w - inner_m * 2
                inner_h = cell_h - inner_m * 2
                scale = min(inner_w / pw, inner_h / ph)
                iw, ih = pw * scale, ph * scale
                ix = cell_x + inner_m + (inner_w - iw) / 2
                iy = cell_y + inner_m + (inner_h - ih) / 2

                buf = io.BytesIO()
                pim.save(buf, "JPEG", quality=90)
                page.insert_image(fitz.Rect(ix, iy, ix + iw, iy + ih), stream=buf.getvalue())
                pim.close()
                img_idx += 1

            # 如果是2合1且有偶数slot用完后画裁剪虚线
            if pages_per_sheet == 2 and img_idx <= total_images and slot == 1:
                _draw_dash_line(page, (0, A4_H / 2), (A4_W, A4_H / 2))

    else:
        doc.close()
        return json_err("支持 1/2/4/6/9 张每页排版")

    out_path = _output_path(fid)
    doc.save(out_path, garbage=4, deflate=True)
    doc.close()
    return send_file(out_path, as_attachment=True,
                   download_name=f"image_print_{pages_per_sheet}up.pdf")


def _draw_page_scaled(target_page, src_page, target_rect, margin=8):
    """
    将源页面渲染到目标矩形内，保持宽高比居中
    """
    inner_rect = fitz.Rect(
        target_rect.x0 + margin, target_rect.y0 + margin,
        target_rect.x1 - margin, target_rect.y1 - margin
    )
    tw = inner_rect.width
    th = inner_rect.height

    # 源页面尺寸
    src_rect = src_page.rect
    sw = src_rect.width
    sh = src_rect.height

    # 计算等比缩放后的尺寸
    scale = min(tw / sw, th / sh)
    img_w = sw * scale
    img_h = sh * scale

    # 居中
    cx = inner_rect.x0 + (tw - img_w) / 2
    cy = inner_rect.y0 + (th - img_h) / 2

    # 渲染
    dpi = max(150, int(scale * 300))
    pix = src_page.get_pixmap(dpi=min(dpi, 300))
    target_page.insert_image(
        fitz.Rect(cx, cy, cx + img_w, cy + img_h),
        stream=pix.tobytes("jpeg")
    )


def _ocr_image_bytes(img_bytes):
    """尝试用Tesseract OCR识别图片文本，失败则返回空字符串"""
    try:
        import pytesseract
        from PIL import Image
        import io
        img = Image.open(io.BytesIO(img_bytes))
        text = pytesseract.image_to_string(img, lang='chi_sim+eng')
        return text
    except ImportError:
        return ""
    except Exception:
        return ""


def _chinese_amount_to_number(chinese_str):
    """将中文大写金额转换为数字

    支持格式如：
    - 壹万贰仟叁佰肆拾伍圆陆角柒分 → 12345.67
    - 壹仟零陆拾元整 → 1060.00
    - 零元整 → 0.00
    - 叁佰圆 → 300.00
    - 壹拾万零伍仟元整 → 105000.00
    - 壹亿贰仟叁佰万圆整 → 123000000.00

    Returns: float 或 None（解析失败）
    """
    import re

    if not chinese_str or not chinese_str.strip():
        return None

    s = chinese_str.strip()

    # 去掉"整"、"正"结尾
    s = re.sub(r'[整正]$', '', s)

    # 全角→半角
    s = s.replace('．', '.').replace('－', '-')

    # 数字映射
    digit_map = {
        '零': 0, '壹': 1, '贰': 2, '叁': 3, '肆': 4,
        '伍': 5, '陆': 6, '柒': 7, '捌': 8, '玖': 9,
        '一': 1, '二': 2, '三': 3, '四': 4, '五': 5,
        '六': 6, '七': 7, '八': 8, '九': 9,
        '０': 0, '１': 1, '２': 2, '３': 3, '４': 4,
        '５': 5, '６': 6, '７': 7, '８': 8, '９': 9,
    }

    # 单位映射
    unit_small = {'角': 0.1, '分': 0.01}  # 小数单位
    unit_section = {'万': 10000, '亿': 100000000}  # 节单位
    unit_base = {'拾': 10, '佰': 100, '仟': 1000}  # 基础单位
    unit_yuan = {'圆': 1, '元': 1}  # 元单位

    # ====== 方法1: 逐字解析法（处理中文大写） ======
    try:
        result = 0.0          # 最终结果
        section_val = 0.0     # 当前节（亿/万之间）的累计值
        current_digit = None  # 当前数字

        i = 0
        while i < len(s):
            ch = s[i]

            # 跳过空格等
            if ch in ' \t':
                i += 1
                continue

            # 数字字符
            if ch in digit_map:
                current_digit = digit_map[ch]
                i += 1
                continue

            # 基础单位（拾佰仟）：当前数字×单位值，累加到节
            if ch in unit_base:
                if current_digit is None:
                    # "拾"前面没有数字，默认为1（如"拾万"=10万）
                    current_digit = 1
                section_val += current_digit * unit_base[ch]
                current_digit = None
                i += 1
                continue

            # 元单位：当前数字直接作为元的个位，节值结算
            if ch in unit_yuan:
                if current_digit is not None:
                    section_val += current_digit
                    current_digit = None
                result += section_val
                section_val = 0.0
                i += 1
                # 检查后面是否还有角/分
                remaining = s[i:]
                decimal_part = 0.0
                j = 0
                while j < len(remaining):
                    rc = remaining[j]
                    if rc in ' \t':
                        j += 1
                        continue
                    if rc in digit_map:
                        # 看下一个是否是角/分
                        if j + 1 < len(remaining) and remaining[j + 1] in unit_small:
                            decimal_part += digit_map[rc] * unit_small[remaining[j + 1]]
                            j += 2
                            continue
                        else:
                            # 数字后面没有角分单位，可能是"零角X分"
                            j += 1
                            continue
                    if rc in unit_small:
                        # 单独出现角/分（前面没有数字），默认为0
                        j += 1
                        continue
                    break
                result += decimal_part
                return result

            # 节单位（万/亿）：节值×节单位，加到总值
            if ch in unit_section:
                if current_digit is not None:
                    section_val += current_digit
                    current_digit = None
                section_val *= unit_section[ch]
                result += section_val
                section_val = 0.0
                i += 1
                continue

            # 角/分（没经过元单位，直接出现角分）
            if ch in unit_small:
                if current_digit is not None:
                    result += current_digit * unit_small[ch]
                    current_digit = None
                i += 1
                continue

            # 未知字符，跳过
            i += 1

        # 如果最后还有未结算的数字和节值
        if current_digit is not None:
            section_val += current_digit
        result += section_val

        if result > 0:
            return result
    except Exception:
        pass

    # ====== 方法2: 简易正则提取（兜底） ======
    # 尝试从字符串中提取纯数字部分
    m = re.search(r'([\d,，.]+)', s)
    if m:
        try:
            return float(m.group(1).replace(',', '').replace('，', ''))
        except ValueError:
            pass

    return None


def _extract_uppercase_total(full_text, form_fields=None):
    """从发票文本中提取大写金额和小写金额

    返回: {
        'uppercase_text': '壹仟零陆拾圆整',    # 大写原文
        'uppercase_value': 1060.0,             # 大写转换后的数字
        'lowercase_text': '¥1060.00',          # 小写原文
        'lowercase_value': 1060.0,             # 小写转换后的数字
    }
    """
    import re

    result = {
        'uppercase_text': '',
        'uppercase_value': None,
        'lowercase_text': '',
        'lowercase_value': None,
    }

    if form_fields is None:
        form_fields = {}

    # ====== 1. 从AcroForm表单域提取大写/小写合计 ======
    # 大写金额字段名
    uppercase_field_names = [
        'jshj_dx', 'jshjDx', 'jshj_daxie', 'jshjDaxie',
        'totalUpper', 'total_upper', 'totalChinese', 'total_chinese',
        'amountChinese', 'amount_chinese', 'amountUpper', 'amount_upper',
        'jshjdx', 'totalAmountChinese', 'totalInWords', 'total_in_words',
        'bigAmount', 'big_amount', 'cnTotal', 'cn_amount',
        'jshj_cn', 'jshjCn', 'amountCn', 'amount_cn',
        'totalCn', 'total_cn', 'totalCnAmount',
    ]
    # 小写金额字段名
    lowercase_field_names = [
        'jshj_xx', 'jshjXx', 'jshj_xiaoxie', 'jshjXiaoxie',
        'totalLower', 'total_lower', 'totalFigure', 'total_figure',
        'amountFigure', 'amount_figure', 'amountNum', 'amount_num',
        'jshjxx', 'totalAmountFigure', 'totalInFigures', 'total_in_figures',
        'smallAmount', 'small_amount', 'jshj_num', 'jshjNum',
        'jshj', 'jshj_hj', 'total', 'totalPrice', 'total_price',
        'amountWithTax', 'amount_with_tax', 'jshjTotal',
    ]

    for fn in uppercase_field_names:
        val = form_fields.get(fn, "")
        if val and val.strip():
            result['uppercase_text'] = val.strip()
            break

    for fn in lowercase_field_names:
        val = form_fields.get(fn, "")
        if val and val.strip():
            result['lowercase_text'] = val.strip()
            break

    # ====== 2. 从文本中提取大写/小写金额 ======
    clean_text = full_text
    compact = re.sub(r'[\s·]+', '', full_text)

    # 提取"价税合计"区域的大写+小写
    # 常见格式:
    #   价税合计（大写）壹仟零陆拾圆整（小写）¥1060.00
    #   价税合计  大写  壹仟零陆拾圆整  小写  ¥1060.00
    #   价税合计（大写金额）壹仟零陆拾圆整  ¥1060.00

    # 模式1: 价税合计...大写...小写
    uppercase_chars = '零壹贰叁肆伍陆柒捌玖拾佰仟万亿圆元角分整正一二三四五六七八九'
    uc_pattern = f'[{uppercase_chars}]+'

    # 先尝试找到"价税合计"区域（匹配到行尾）
    total_area_match = re.search(
        r'价\s*税\s*合\s*计[^\n]*', clean_text
    )
    if total_area_match:
        total_area = total_area_match.group(0)
        # 在这个区域中找大写金额
        if not result['uppercase_text']:
            # 大写金额：在"大写"后面，到"小写"或"¥"或换行之前
            m = re.search(
                r'(?:大写[金额）)\s:]*)\s*(' + uc_pattern + r')',
                total_area
            )
            if m:
                result['uppercase_text'] = m.group(1)
            else:
                # 没有明确的"大写"标记，尝试找连续中文数字串
                m = re.search(r'(' + uc_pattern + r'[圆元][整正]?)', total_area)
                if m:
                    candidate = m.group(1)
                    # 必须含有单位才认定
                    if any(u in candidate for u in '拾佰仟万亿圆元'):
                        result['uppercase_text'] = candidate

        if not result['lowercase_text']:
            # 小写金额：¥后面的数字（优先匹配有小数点的，避免误匹配发票号码）
            m = re.search(r'[¥￥]\s*([\d,，]+\.\d{1,2})', total_area)
            if m:
                result['lowercase_text'] = '¥' + m.group(1)
            else:
                # 没有小数点的金额，验证不是8位以上纯整数（发票号码）
                m = re.search(r'[¥￥]\s*([\d,，]+)', total_area)
                if m:
                    num_str = m.group(1).replace(',', '').replace('，', '')
                    if not re.match(r'^\d{8,}$', num_str):  # 排除发票号码
                        result['lowercase_text'] = '¥' + m.group(1)

    # 如果上面没找到，扩大搜索范围
    if not result['uppercase_text']:
        # 全文搜索大写金额
        m = re.search(r'(?:大写[金额）)\s:]*)\s*(' + uc_pattern + r')', clean_text)
        if m:
            result['uppercase_text'] = m.group(1)

    if not result['lowercase_text']:
        # 全文搜索价税合计后面的小写金额（优先有小数点）
        m = re.search(r'价\s*税\s*合\s*计[^¥￥\n]*[¥￥]\s*([\d,，]+\.\d{1,2})', clean_text)
        if m:
            result['lowercase_text'] = '¥' + m.group(1)
        else:
            m = re.search(r'价\s*税\s*合\s*计[^¥￥\n]*[¥￥]\s*([\d,，]+)', clean_text)
            if m:
                num_str = m.group(1).replace(',', '').replace('，', '')
                if not re.match(r'^\d{8,}$', num_str):  # 排除发票号码
                    result['lowercase_text'] = '¥' + m.group(1)

    # ====== 3. 转换大写→数字 ======
    if result['uppercase_text']:
        converted = _chinese_amount_to_number(result['uppercase_text'])
        if converted is not None:
            result['uppercase_value'] = round(converted, 2)

    # ====== 4. 转换小写→数字 ======
    if result['lowercase_text']:
        m = re.search(r'([\d,，.]+)', result['lowercase_text'])
        if m:
            try:
                result['lowercase_value'] = round(float(m.group(1).replace(',', '').replace('，', '')), 2)
            except ValueError:
                pass

    return result


def _extract_invoice_text(fpath):
    """从发票PDF中提取文本，使用三种策略：文本流→AcroForm表单域→OCR图片

    返回 (full_text, form_fields, has_ocr)
    - full_text: 所有文本拼接
    - form_fields: {field_name: field_value} 从AcroForm提取
    - has_ocr: 是否使用了OCR降级
    """
    doc = fitz.open(fpath)
    all_text = []
    form_fields = {}
    has_ocr = False

    for page_num in range(doc.page_count):
        pg = doc[page_num]

        # 策略1: AcroForm表单域提取（全电发票/电子发票主要方式）
        try:
            for widget in pg.widgets():
                fname = widget.field_name or ""
                fvalue = widget.field_value or ""
                if fname and fvalue:
                    form_fields[fname] = fvalue
        except Exception:
            pass

        # 策略2: 常规文本提取
        text = pg.get_text("text")
        text_quality = len(text.replace("·", "").replace(" ", "").strip())

        if text_quality < 20:
            # 文本太少或全是点号，尝试dict模式获取更多细节
            try:
                d = pg.get_text("dict")
                dict_text_parts = []
                for block in d.get("blocks", []):
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line["spans"]:
                                t = span.get("text", "").strip()
                                if t:
                                    dict_text_parts.append(t)
                dict_text = "\n".join(dict_text_parts)
                if len(dict_text.replace("·", "").strip()) > text_quality:
                    text = dict_text
            except Exception:
                pass

        # 策略3: 如果文本仍然太少，尝试OCR
        final_text = text.replace("·", "").replace(" ", "")
        if len(final_text.strip()) < 15:
            try:
                pix = pg.get_pixmap(dpi=200)
                img_data = pix.tobytes("png")
                ocr_text = _ocr_image_bytes(img_data)
                if ocr_text and len(ocr_text.strip()) > len(final_text.strip()):
                    text = ocr_text
                    has_ocr = True
            except Exception:
                pass

        all_text.append(text)

    doc.close()
    full_text = "\n".join(all_text)
    return full_text, form_fields, has_ocr


def _draw_dash_line(page, p1, p2):
    """画虚线"""
    x1, y1 = p1
    x2, y2 = p2
    seg_len = 8
    gap_len = 4
    total_len = seg_len + gap_len
    length = ((x2 - x1) ** 2 + (y2 - y1) ** 2) ** 0.5
    if length == 0:
        return
    segments = int(length / total_len)
    dx = (x2 - x1) / length
    dy = (y2 - y1) / length
    for s in range(segments):
        sx = x1 + dx * s * total_len
        sy = y1 + dy * s * total_len
        ex = x1 + dx * (s * total_len + seg_len)
        ey = y1 + dy * (s * total_len + seg_len)
        page.draw_line((sx, sy), (ex, ey), color=(0.6, 0.6, 0.6), width=0.5)


def _parse_invoice_text(full_text, form_fields=None):
    """从发票文本和表单域中提取结构化信息，支持多种发票格式

    Args:
        full_text: PDF提取的文本
        form_fields: AcroForm表单域字典 {field_name: field_value}

    Returns:
        dict: 结构化发票信息
    """
    import re

    if form_fields is None:
        form_fields = {}

    info = {
        "invoice_type": "",
        "invoice_code": "",
        "invoice_number": "",
        "invoice_date": "",
        "seller": "",
        "buyer": "",
        "amount": "",
        "tax": "",
        "total": "",
        "total_uppercase": "",       # 大写金额原文
        "total_lowercase": "",       # 小写金额原文
        "total_verified": False,     # 大写/小写是否一致
        "total_verify_msg": "",      # 验证说明
    }

    # ============ 1. 从AcroForm表单域提取（优先级最高） ============
    # 常见发票表单域字段名映射
    field_map = {
        # 发票类型
        "invoice_type": ["fptype", "fp_type", "fpType", "fplxdm", "invoiceType",
                        "fpzl", "invoiceKind", "type"],
        # 发票代码
        "invoice_code": ["fpdm", "fp_dm", "fpDm", "invoiceCode", "invoice_code",
                        "invoiceCode1", "invoice_code1"],
        # 发票号码
        "invoice_number": ["fphm", "fp_hm", "fpHm", "invoiceNumber", "invoice_number",
                          "invoiceNo", "invoice_no", "fpdmhm", "einvNum"],
        # 开票日期
        "invoice_date": ["kprq", "kp_rq", "kpRq", "invoiceDate", "invoice_date",
                        "kprqq", "billingDate", "billing_date", "makeDate",
                        "issueDate", "kprq_formatted"],
        # 销售方名称
        "seller": ["xfmc", "xf_mc", "xfMc", "sellerName", "seller_name",
                  "xfname", "saleName", "sale_name", "saleTaxName",
                  "sellName", "payeeName", "xfmc2"],
        # 购买方名称
        "buyer": ["gfmc", "gf_mc", "gfMc", "buyerName", "buyer_name",
                 "gfname", "purchaseName", "purchase_name", "payerName",
                 "buyName", "gfmc2"],
        # 金额
        "amount": ["je", "jshj_xe", "amount", "hjje", "totalAmount",
                  "total_amount", "sumAmount", "jshj_je", "noTaxAmount"],
        # 税额
        "tax": ["se", "tax", "hjse", "totalTax", "total_tax",
               "sumTax", "jshj_se"],
        # 价税合计（小写）
        "total": ["jshj", "jshj_hj", "total", "totalPrice", "total_price",
                 "amountWithTax", "amount_with_tax", "jshjTotal",
                 "jshj_xx", "jshjXx", "jshj_xiaoxie", "jshjXiaoxie",
                 "totalLower", "total_lower", "totalFigure", "total_figure",
                 "jshjxx", "totalInFigures", "total_in_figures",
                 "jshj_num", "jshjNum", "jshjJe"],
        # 价税合计（大写原文）
        "total_uppercase": ["jshj_dx", "jshjDx", "jshj_daxie", "jshjDaxie",
                           "totalUpper", "total_upper", "totalChinese", "total_chinese",
                           "amountChinese", "amount_chinese", "amountUpper", "amount_upper",
                           "jshjdx", "totalAmountChinese", "totalInWords", "total_in_words",
                           "bigAmount", "big_amount", "cnTotal", "cn_amount",
                           "jshj_cn", "jshjCn", "amountCn", "amount_cn",
                           "totalCn", "total_cn", "totalCnAmount", "jshjDxJe",
                           "jshj_hz", "jshjHz"],
    }

    for target, field_names in field_map.items():
        for fn in field_names:
            val = form_fields.get(fn, "")
            if val and val.strip():
                info[target] = val.strip()
                break

    # 格式化日期（从"2024年06月15日"→"2024-06-15"）
    if info["invoice_date"]:
        dm = re.match(r"(\d{4})\s*年\s*(\d{1,2})\s*月\s*(\d{1,2})\s*日?", info["invoice_date"])
        if dm:
            info["invoice_date"] = f"{dm.group(1)}-{dm.group(2).zfill(2)}-{dm.group(3).zfill(2)}"
        else:
            dm = re.match(r"(\d{4})[/\-](\d{1,2})[/\-](\d{1,2})", info["invoice_date"])
            if dm:
                info["invoice_date"] = f"{dm.group(1)}-{dm.group(2).zfill(2)}-{dm.group(3).zfill(2)}"

    # ============ 2. 从文本中补充提取（表单域可能不全） ============
    # 合并文本（去点号增强匹配）
    clean_text = full_text
    # 也生成一个"紧凑版"文本（去掉所有空白和中间点，便于匹配）
    compact = re.sub(r'[\s·]+', '', full_text)

    # 发票类型（如果表单域没有）
    if not info["invoice_type"]:
        type_patterns = [
            ("增值税电子普通发票", "增值税电子普通发票"),
            ("增值税电子专用发票", "增值税电子专用发票"),
            ("增值税普通发票", "增值税普通发票"),
            ("增值税专用发票", "增值税专用发票"),
            ("电子发票（普通发票）", "电子发票（普通发票）"),
            ("电子发票（增值税普通发票）", "电子发票（增值税普通发票）"),
            ("电子发票（增值税专用发票）", "电子发票（增值税专用发票）"),
            ("全电发票", "全电发票"),
            ("电子发票", "电子发票"),
            ("机打发票", "机打发票"),
            ("机动车销售统一发票", "机动车销售统一发票"),
            ("二手车销售统一发票", "二手车销售统一发票"),
        ]
        for pattern, label in type_patterns:
            if pattern in clean_text or pattern in compact:
                info["invoice_type"] = label
                break
        if not info["invoice_type"]:
            info["invoice_type"] = "其他"

    # 发票代码（如果表单域没有）
    if not info["invoice_code"]:
        code_patterns = [
            r"发票代码[：:\s·]*(\d{10,12})",
            r"发\s*票\s*代\s*码[：:\s·]*(\d{10,12})",
        ]
        for pat in code_patterns:
            m = re.search(pat, clean_text)
            if m:
                info["invoice_code"] = m.group(1)
                break

    # 发票号码（如果表单域没有）
    if not info["invoice_number"]:
        number_patterns = [
            r"发票号码[：:\s·]*(\d{8,20})",
            r"发\s*票\s*号\s*码[：:\s·]*(\d{8,20})",
            r"No[.：:\s]*(\d{8,20})",
        ]
        for pat in number_patterns:
            m = re.search(pat, clean_text)
            if m:
                info["invoice_number"] = m.group(1)
                break

    # 开票日期（如果表单域没有）
    if not info["invoice_date"]:
        date_patterns = [
            r"(?:开票日期|开具日期)[：:\s]*(\d{4})\s*年\s*(\d{1,2})\s*月\s*(\d{1,2})\s*日",
            r"(?:开票日期|开具日期)[：:\s]*(\d{4})[年/\-·\s](\d{1,2})[月/\-·\s](\d{1,2})",
            r"(\d{4})\s*年\s*(\d{1,2})\s*月\s*(\d{1,2})\s*日",
        ]
        for pat in date_patterns:
            m = re.search(pat, clean_text)
            if m:
                info["invoice_date"] = f"{m.group(1)}-{m.group(2).zfill(2)}-{m.group(3).zfill(2)}"
                break

    # 销售方名称（如果表单域没有）
    if not info["seller"]:
        seller_patterns = [
            r"(?:销售方|销方|收款单位|卖方)[名：:\s]*称[：:\s]*([^\n,，\s]{2,50})",
            r"(?:销售方|销方|收款单位|卖方)[名称：:\s]*([^\n,，\s]{2,50})",
            r"名\s*称[：:\s]*([^\n,，\s]{2,50})\s*(?:纳税人识别号|统一社会信用代码)",
        ]
        for pat in seller_patterns:
            m = re.search(pat, clean_text)
            if m:
                info["seller"] = m.group(1).strip()
                break

    # 购买方名称（如果表单域没有）
    if not info["buyer"]:
        buyer_patterns = [
            r"(?:购买方|购方|付款单位|买方)[名：:\s]*称[：:\s]*([^\n,，\s]{2,50})",
            r"(?:购买方|购方|付款单位|买方)[名称：:\s]*([^\n,，\s]{2,50})",
        ]
        for pat in buyer_patterns:
            m = re.search(pat, clean_text)
            if m:
                info["buyer"] = m.group(1).strip()
                break

    # 金额（如果表单域没有）
    # 注意：金额提取必须严格上下文约束，避免把发票号码等数字误识别为金额
    if not info["amount"]:
        amount_patterns = [
            # 明确的"合计金额"或"金额"关键字后面跟¥符号
            r"合\s*计\s*金\s*额[：:\s]*[¥￥]\s*([\d,，]+\.\d{1,2})",
            r"金\s*额[：:\s]*[¥￥]\s*([\d,，]+\.\d{1,2})",
            r"合\s*计[：:\s]*[¥￥]\s*([\d,，]+\.\d{1,2})",
            r"Amount[：:\s]*[¥￥]\s*([\d,，]+\.\d{1,2})",
            # 没有¥符号，但紧跟"金额"关键字且有小数点的数字
            r"金\s*额[：:\s]*(\d+\.\d{1,2})",
            r"合\s*计\s*金\s*额[：:\s]*(\d+\.\d{1,2})",
        ]
        for pat in amount_patterns:
            m = re.search(pat, clean_text)
            if m:
                info["amount"] = m.group(1)
                break
        # 如果还是没有，尝试在紧凑文本中查找（仍然要求有小数点）
        if not info["amount"]:
            m = re.search(r"金\s*额[：:\s]*(\d+\.\d{1,2})", compact)
            if m:
                info["amount"] = m.group(1)

    # 税额（如果表单域没有）
    if not info["tax"]:
        tax_patterns = [
            r"合\s*计\s*税\s*额[：:\s]*[¥￥]\s*([\d,，]+\.\d{1,2})",
            r"税\s*额[：:\s]*[¥￥]\s*([\d,，]+\.\d{1,2})",
            r"税\s*额[：:\s]*(\d+\.\d{1,2})",
            r"Tax[：:\s]*[¥￥]?\s*([\d,，]+\.\d{1,2})",
        ]
        for pat in tax_patterns:
            m = re.search(pat, clean_text)
            if m:
                info["tax"] = m.group(1)
                break

    # 价税合计（如果表单域没有）
    if not info["total"]:
        total_patterns = [
            # 价税合计 + ¥/￥ + 小数金额
            r"价\s*税\s*合\s*计[^¥￥\n]*[¥￥]\s*([\d,，]+\.\d{1,2})",
            # 价税合计后面紧跟数字（带小数点）
            r"价\s*税\s*合\s*计[：:\s]*[¥￥]?\s*(\d+\.\d{1,2})",
        ]
        for pat in total_patterns:
            m = re.search(pat, clean_text)
            if m:
                info["total"] = m.group(1)
                break
        # 如果还是没有，尝试在紧凑文本中查找（要求小数点）
        if not info["total"]:
            m = re.search(r"价\s*税\s*合\s*计[^\d]*([\d,，]+\.\d{1,2})", compact)
            if m:
                info["total"] = m.group(1)

    # ============ 2.5 大写/小写金额交叉验证 ============
    # 调用 _extract_uppercase_total 从文本和表单域中提取大写和小写金额
    uc_lc = _extract_uppercase_total(full_text, form_fields)

    # 如果表单域已经提取到大写金额但 _extract_uppercase_total 没有找到，使用表单域的值
    if info["total_uppercase"] and not uc_lc['uppercase_text']:
        uc_lc['uppercase_text'] = info["total_uppercase"]
        # 同时转换
        converted = _chinese_amount_to_number(info["total_uppercase"])
        if converted is not None:
            uc_lc['uppercase_value'] = round(converted, 2)

    # 如果已有 total 但 _extract_uppercase_total 没找到小写，用已有的 total
    if info["total"] and not uc_lc['lowercase_text']:
        uc_lc['lowercase_text'] = '¥' + str(info["total"])
        try:
            uc_lc['lowercase_value'] = round(float(str(info["total"]).replace(",", "").replace("，", "")), 2)
        except ValueError:
            pass

    # 如果 _extract_uppercase_total 找到了小写但 total 还是空，使用它
    if not info["total"] and uc_lc['lowercase_value'] is not None:
        info["total"] = str(uc_lc['lowercase_value'])

    # 如果表单域没有大写但 _extract_uppercase_total 找到了，也保存
    if not info["total_uppercase"] and uc_lc['uppercase_text']:
        info["total_uppercase"] = uc_lc['uppercase_text']

    # 保存小写原文
    if uc_lc['lowercase_text']:
        info["total_lowercase"] = uc_lc['lowercase_text']

    # ============ 执行交叉验证 ============
    uc_val = uc_lc['uppercase_value']    # 大写转数字
    lc_val = uc_lc['lowercase_value']    # 小写数字

    if uc_val is not None and lc_val is not None:
        # 大写和小写都有值，进行交叉比对
        if abs(uc_val - lc_val) < 0.01:
            # 一致：使用小写值（更精确），标记为已验证
            info["total"] = str(lc_val)
            info["total_verified"] = True
            info["total_verify_msg"] = f"✅ 大写({uc_val:.2f})与小写({lc_val:.2f})一致"
        else:
            # 不一致：优先使用大写金额（大写通常更可靠），标记不一致
            info["total"] = str(uc_val)
            info["total_verified"] = False
            info["total_verify_msg"] = f"⚠️ 大写({uc_val:.2f})与小写({lc_val:.2f})不一致，已采用大写值"
    elif uc_val is not None:
        # 只有大写，转换后使用
        info["total"] = str(uc_val)
        info["total_verified"] = False
        info["total_verify_msg"] = f"仅有大写金额({uc_val:.2f})，无法交叉验证"
    elif lc_val is not None:
        # 只有小写，直接使用
        info["total"] = str(lc_val)
        info["total_verified"] = False
        info["total_verify_msg"] = f"仅有小写金额({lc_val:.2f})，无法交叉验证"
    else:
        info["total_verify_msg"] = "未找到大写或小写金额"

    # ============ 3. 从表单域值直接匹配文本中缺失的字段 ============
    # 有些表单域的field_name不标准，但field_value可能包含关键字
    # 重要：需要区分发票号码(8-20位整数)、发票代码(10-12位整数)、金额(通常有小数)
    if form_fields:
        # 收集所有未匹配到映射的纯数字/日期字段值，分类处理
        unmatched_numeric = []  # (field_name, value, type_hint)
        for fname, fvalue in form_fields.items():
            fv = fvalue.strip()
            if not fv:
                continue

            # 如果value包含"发票"字样但类型还没识别到
            if "发票" in fv and not info["invoice_type"]:
                for tt in ["增值税电子普通发票", "增值税电子专用发票", "增值税普通发票",
                          "增值税专用发票", "电子发票", "全电发票", "机打发票"]:
                    if tt in fv:
                        info["invoice_type"] = tt
                        break

            # 如果value是日期格式但日期还没识别
            if not info["invoice_date"]:
                dm = re.match(r"(\d{4})[年/\-](\d{1,2})[月/\-](\d{1,2})", fv)
                if dm:
                    info["invoice_date"] = f"{dm.group(1)}-{dm.group(2).zfill(2)}-{dm.group(3).zfill(2)}"
                    continue

            # 智能数字分类
            # 发票代码：10-12位纯整数
            if re.match(r"^\d{10,12}$", fv) and not info["invoice_code"]:
                info["invoice_code"] = fv
                continue

            # 发票号码：8位纯整数（传统发票），或更长的纯整数（全电发票20位）
            if re.match(r"^\d{8}$", fv) and not info["invoice_number"]:
                info["invoice_number"] = fv
                continue
            if re.match(r"^\d{10,20}$", fv) and not info["invoice_number"]:
                # 长整数：可能是发票号码（全电发票号码20位），不是金额
                # 金额极少是10位以上的纯整数（超过99亿的概率极低）
                info["invoice_number"] = fv
                continue

            # 金额特征：包含小数点，或数值合理（< 10亿），或有逗号分隔
            if re.match(r"^[\d,，.]+$", fv):
                try:
                    num_val = float(fv.replace(",", "").replace("，", ""))
                except ValueError:
                    num_val = None

                if num_val is not None and num_val > 0:
                    # 关键判断：纯整数 >= 8位 → 大概率是发票号码/代码，不是金额
                    pure_int = re.match(r"^\d+$", fv)
                    if pure_int and len(fv) >= 8:
                        # 8位以上纯整数，优先归类为发票号码
                        if not info["invoice_number"]:
                            info["invoice_number"] = fv
                        # 不要当成金额！
                    elif "." in fv or num_val < 100000000:
                        # 有小数点或数值合理(<1亿) → 可能是金额
                        # 但需要根据字段名判断应该放到 amount / tax / total 哪个
                        fname_lower = fname.lower()
                        if any(k in fname_lower for k in ['se', 'tax', 'taxamount', 'hjse']):
                            if not info["tax"]:
                                info["tax"] = fv
                        elif any(k in fname_lower for k in ['je', 'amount', 'hjje', 'noTax', 'xiaoxie', 'xx', 'num', 'figure']):
                            if not info["amount"]:
                                info["amount"] = fv
                        elif any(k in fname_lower for k in ['jshj', 'total', 'hj', 'price', 'withtax']):
                            if not info["total"]:
                                info["total"] = fv
                        elif not info["total"]:
                            # 字段名没有明确提示，放到total（但仅当值看起来像金额）
                            if "." in fv or num_val < 100000:
                                info["total"] = fv

    # ============ 4. 从OCR文本中提取（如果有OCR标记） ============
    # 如果金额仍然为空，尝试从整个文本中找带¥/￥的金额
    # 重要：必须校验数字看起来像金额（有小数点），避免把发票号码当金额
    if not info["total"]:
        # 优先找带小数点的金额（真正的金额格式）
        all_amounts = re.findall(r"[¥￥]\s*([\d,，]+\.\d{1,2})", clean_text)
        if all_amounts:
            # 取最后一个带小数点的（通常是价税合计）
            info["total"] = all_amounts[-1]
        else:
            # 没有小数点的¥金额，需要谨慎：验证不是发票号码
            # 发票号码通常是8-20位纯整数，金额通常<10亿
            all_num_amounts = re.findall(r"[¥￥]\s*([\d,，]+)", clean_text)
            for amt in reversed(all_num_amounts):  # 从后往前找
                try:
                    num_val = float(amt.replace(",", "").replace("，", ""))
                    # 排除明显是发票号码的值：8位以上纯整数
                    if re.match(r"^\d{8,}$", amt.strip()):
                        continue  # 跳过，很可能是发票号码
                    if 0 < num_val < 100000000:  # 合理金额范围（<1亿）
                        info["total"] = amt
                        break
                except ValueError:
                    continue

    if not info["amount"] and info["total"]:
        # 如果只有total没有amount，尝试从total反推
        try:
            total_val = float(str(info["total"]).replace(",", "").replace("，", ""))
            tax_val = float(str(info["tax"]).replace(",", "").replace("，", "")) if info["tax"] else 0
            if total_val > 0 and tax_val >= 0:
                info["amount"] = f"{total_val - tax_val:.2f}"
        except ValueError:
            pass

    return info


# ============================================================
# 辅助函数
# ============================================================

def _find_upload(fid):
    """根据ID查找上传文件"""
    for f in UPLOAD_DIR.iterdir():
        if f.stem == fid:
            return str(f)
    return None


def _output_path(fid, ext=".pdf"):
    p = OUTPUT_DIR / f"{fid}{ext}"
    p.parent.mkdir(exist_ok=True)
    return str(p)


def _parse_color(hex_color):
    """#RRGGBB -> (r, g, b) 0-1"""
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return (r / 255, g / 255, b / 255)


def _parse_color_rgb(hex_color):
    """#RRGGBB -> (r, g, b) 0-255"""
    hex_color = hex_color.lstrip("#")
    return (int(hex_color[0:2], 16) / 255, int(hex_color[2:4], 16) / 255, int(hex_color[4:6], 16) / 255)


# ============================================================
# 15. 发票功能
# ============================================================

@app.route("/api/invoice/print", methods=["POST"])
def invoice_print():
    """发票打印：将多张发票PDF排版到A4纸上，支持1/2/4张每页"""
    per_sheet = int(request.form.get("per_sheet", 1))
    files = request.files.getlist("files")
    if not files:
        f = request.files.get("file")
        if not f:
            return json_err("请上传发票文件")
        files = [f]

    try:
        # 先合并所有发票
        doc = fitz.open()
        fid = None
        for f in files:
            fpath, _fid = save_upload(f)
            fid = _fid
            src = fitz.open(fpath)
            doc.insert_pdf(src)
            src.close()

        if per_sheet == 1:
            # 每页1张发票，直接输出
            out_path = _output_path(fid)
            doc.save(out_path, garbage=4, deflate=True)
            doc.close()
            return send_file(out_path, as_attachment=True, download_name="invoice_print.pdf")

        # 多合一排版
        doc_out = fitz.open()
        A4_W, A4_H = 595, 842

        if per_sheet == 2:
            # 上下布局，含裁剪线
            gap = 6
            line_y = A4_H / 2
            top_h = line_y - gap / 2
            bot_h = A4_H - line_y - gap / 2

            for i in range(0, doc.page_count, 2):
                page_a4 = doc_out.new_page(width=A4_W, height=A4_H)
                # 上半部分
                _draw_page_scaled(page_a4, doc[i],
                    fitz.Rect(0, 0, A4_W, top_h), margin=8)
                # 裁剪虚线
                _draw_dash_line(page_a4, (0, line_y), (A4_W, line_y))
                # 下半部分
                if i + 1 < doc.page_count:
                    _draw_page_scaled(page_a4, doc[i + 1],
                        fitz.Rect(0, line_y + gap / 2, A4_W, A4_H), margin=8)

        elif per_sheet == 4:
            half_w = A4_W / 2
            half_h = A4_H / 2
            margin = 4
            positions = [
                fitz.Rect(margin, margin, half_w - margin, half_h - margin),
                fitz.Rect(half_w + margin, margin, A4_W - margin, half_h - margin),
                fitz.Rect(margin, half_h + margin, half_w - margin, A4_H - margin),
                fitz.Rect(half_w + margin, half_h + margin, A4_W - margin, A4_H - margin),
            ]
            for i in range(0, doc.page_count, 4):
                page_a4 = doc_out.new_page(width=A4_W, height=A4_H)
                for j, rect in enumerate(positions):
                    if i + j < doc.page_count:
                        _draw_page_scaled(page_a4, doc[i + j], rect, margin=0)
                # 画十字裁剪线
                _draw_dash_line(page_a4, (0, A4_H / 2), (A4_W, A4_H / 2))
                _draw_dash_line(page_a4, (A4_W / 2, 0), (A4_W / 2, A4_H))

        out_path = _output_path(fid)
        doc_out.save(out_path)
        doc_out.close()
        doc.close()
        return send_file(out_path, as_attachment=True,
                       download_name=f"invoice_{per_sheet}up.pdf")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/invoice/merge", methods=["POST"])
def invoice_merge():
    """发票合并：将多张发票PDF合并为一个文件"""
    files = request.files.getlist("files")
    if not files:
        return json_err("请上传发票文件")

    try:
        doc_out = fitz.open()
        fid = None
        for f in files:
            fpath, _fid = save_upload(f)
            fid = _fid
            src = fitz.open(fpath)
            doc_out.insert_pdf(src)
            src.close()

        out_path = _output_path(fid)
        doc_out.save(out_path, garbage=4, deflate=True)
        doc_out.close()
        return send_file(out_path, as_attachment=True,
                       download_name="invoices_merged.pdf")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/invoice/organize", methods=["POST"])
def invoice_organize():
    """发票整理：将多张发票排版到A4纸上，每页2张含裁剪线，适合裁剪归档"""
    files = request.files.getlist("files")
    if not files:
        return json_err("请上传发票文件")

    per_sheet = int(request.form.get("per_sheet", 2))

    try:
        doc_all = fitz.open()
        fid = None
        for f in files:
            fpath, _fid = save_upload(f)
            fid = _fid
            src = fitz.open(fpath)
            doc_all.insert_pdf(src)
            src.close()

        doc_out = fitz.open()
        A4_W, A4_H = 595, 842

        if per_sheet == 2:
            # 标准发票整理：每页2张，上下布局
            gap = 8
            line_y = A4_H / 2
            top_h = line_y - gap / 2
            bot_h = A4_H - line_y - gap / 2

            for i in range(0, doc_all.page_count, 2):
                page_a4 = doc_out.new_page(width=A4_W, height=A4_H)
                # 页码标注
                page_a4.insert_text((A4_W - 60, 14),
                    f"第 {i // 2 + 1} 页", fontsize=8, color=(0.6, 0.6, 0.6),
                    fontname="hebo", fontfile=None)
                # 上半部分
                _draw_page_scaled(page_a4, doc_all[i],
                    fitz.Rect(0, 20, A4_W, top_h), margin=6)
                # 裁剪虚线（带剪刀标记）
                _draw_dash_line(page_a4, (20, line_y), (A4_W - 20, line_y))
                page_a4.insert_text((6, line_y - 2), "✂",
                    fontsize=10, color=(0.6, 0.6, 0.6), fontname="hebo", fontfile=None)
                # 下半部分
                if i + 1 < doc_all.page_count:
                    _draw_page_scaled(page_a4, doc_all[i + 1],
                        fitz.Rect(0, line_y + gap / 2, A4_W, A4_H), margin=6)

        elif per_sheet == 4:
            # 每页4张，2x2网格
            half_w = A4_W / 2
            half_h = A4_H / 2
            margin = 4
            positions = [
                fitz.Rect(margin, 20 + margin, half_w - margin, half_h - margin),
                fitz.Rect(half_w + margin, 20 + margin, A4_W - margin, half_h - margin),
                fitz.Rect(margin, half_h + margin, half_w - margin, A4_H - margin),
                fitz.Rect(half_w + margin, half_h + margin, A4_W - margin, A4_H - margin),
            ]
            for i in range(0, doc_all.page_count, 4):
                page_a4 = doc_out.new_page(width=A4_W, height=A4_H)
                page_a4.insert_text((A4_W - 60, 14),
                    f"第 {i // 4 + 1} 页", fontsize=8, color=(0.6, 0.6, 0.6),
                    fontname="hebo", fontfile=None)
                for j, rect in enumerate(positions):
                    if i + j < doc_all.page_count:
                        _draw_page_scaled(page_a4, doc_all[i + j], rect, margin=0)
                # 十字裁剪线
                _draw_dash_line(page_a4, (20, A4_H / 2), (A4_W - 20, A4_H / 2))
                _draw_dash_line(page_a4, (A4_W / 2, 20), (A4_W / 2, A4_H - 20))
                page_a4.insert_text((6, A4_H / 2 - 2), "✂",
                    fontsize=10, color=(0.6, 0.6, 0.6), fontname="hebo", fontfile=None)

        out_path = _output_path(fid)
        doc_out.save(out_path, garbage=4, deflate=True)
        doc_out.close()
        doc_all.close()
        return send_file(out_path, as_attachment=True,
                       download_name=f"invoices_organized_{per_sheet}up.pdf")
    except Exception as e:
        return json_err(str(e))


@app.route("/api/invoice/parse", methods=["POST"])
def invoice_parse():
    """发票解析：从发票PDF中提取文本信息，支持AcroForm+OCR"""
    files = request.files.getlist("files")
    if not files:
        f = request.files.get("file")
        if not f:
            return json_err("请上传发票文件")
        files = [f]

    results = []
    try:
        for f in files:
            fpath, fid = save_upload(f)
            full_text, form_fields, has_ocr = _extract_invoice_text(fpath)

            info = _parse_invoice_text(full_text, form_fields)
            info["filename"] = f.filename
            info["fid"] = fid
            info["raw_text"] = full_text[:5000]
            info["form_fields"] = form_fields
            info["ocr_used"] = has_ocr

            results.append(info)

        return json_ok(results)
    except Exception as e:
        return json_err(str(e))


@app.route("/api/invoice/statistics", methods=["POST"])
def invoice_statistics():
    """发票统计：汇总多张发票的总金额、税额、按月/类型分组统计，支持导出Excel"""
    files = request.files.getlist("files")
    if not files:
        return json_err("请上传发票文件")

    try:
        invoices = []
        for f in files:
            fpath, fid = save_upload(f)
            full_text, form_fields, has_ocr = _extract_invoice_text(fpath)
            parsed = _parse_invoice_text(full_text, form_fields)

            def safe_float(s):
                try:
                    return float(str(s).replace(",", "").replace("，", ""))
                except (ValueError, TypeError):
                    return 0.0

            info = {
                "filename": f.filename,
                "fid": fid,
                "invoice_type": parsed.get("invoice_type", "其他"),
                "invoice_code": parsed.get("invoice_code", ""),
                "invoice_number": parsed.get("invoice_number", ""),
                "invoice_date": parsed.get("invoice_date", ""),
                "seller": parsed.get("seller", ""),
                "buyer": parsed.get("buyer", ""),
                "amount": safe_float(parsed.get("amount", 0)),
                "tax": safe_float(parsed.get("tax", 0)),
                "total": safe_float(parsed.get("total", 0)),
                "total_uppercase": parsed.get("total_uppercase", ""),
                "total_lowercase": parsed.get("total_lowercase", ""),
                "total_verified": parsed.get("total_verified", False),
                "total_verify_msg": parsed.get("total_verify_msg", ""),
            }

            # 如果total为0但amount有值，用amount+tax
            if info["total"] == 0 and info["amount"] > 0:
                info["total"] = round(info["amount"] + info["tax"], 2)

            invoices.append(info)

        # 汇总统计
        total_amount = sum(inv["amount"] for inv in invoices)
        total_tax = sum(inv["tax"] for inv in invoices)
        total_total = sum(inv["total"] for inv in invoices)

        # 按月份分组
        monthly = {}
        for inv in invoices:
            month_key = inv["invoice_date"][:7] if inv["invoice_date"] else "未知"
            if month_key not in monthly:
                monthly[month_key] = {"count": 0, "amount": 0.0, "tax": 0.0, "total": 0.0}
            monthly[month_key]["count"] += 1
            monthly[month_key]["amount"] += inv["amount"]
            monthly[month_key]["tax"] += inv["tax"]
            monthly[month_key]["total"] += inv["total"]

        # 按类型分组
        by_type = {}
        for inv in invoices:
            t = inv["invoice_type"] or "其他"
            if t not in by_type:
                by_type[t] = {"count": 0, "amount": 0.0, "tax": 0.0, "total": 0.0}
            by_type[t]["count"] += 1
            by_type[t]["amount"] += inv["amount"]
            by_type[t]["tax"] += inv["tax"]
            by_type[t]["total"] += inv["total"]

        stats = {
            "total_count": len(invoices),
            "total_amount": round(total_amount, 2),
            "total_tax": round(total_tax, 2),
            "total_total": round(total_total, 2),
            "monthly": monthly,
            "by_type": by_type,
            "invoices": invoices,
        }

        return json_ok(stats)
    except Exception as e:
        return json_err(str(e))


@app.route("/api/invoice/export", methods=["POST"])
def invoice_export():
    """发票统计导出Excel：上传发票PDF，识别后导出为Excel统计表"""
    files = request.files.getlist("files")
    if not files:
        return json_err("请上传发票文件")

    try:
        import openpyxl
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side

        invoices = []
        for f in files:
            fpath, fid = save_upload(f)
            full_text, form_fields, has_ocr = _extract_invoice_text(fpath)
            parsed = _parse_invoice_text(full_text, form_fields)

            def safe_float(s):
                try:
                    return float(str(s).replace(",", "").replace("，", ""))
                except (ValueError, TypeError):
                    return 0.0

            info = {
                "filename": f.filename,
                "invoice_type": parsed.get("invoice_type", "其他"),
                "invoice_code": parsed.get("invoice_code", ""),
                "invoice_number": parsed.get("invoice_number", ""),
                "invoice_date": parsed.get("invoice_date", ""),
                "seller": parsed.get("seller", ""),
                "buyer": parsed.get("buyer", ""),
                "amount": safe_float(parsed.get("amount", 0)),
                "tax": safe_float(parsed.get("tax", 0)),
                "total": safe_float(parsed.get("total", 0)),
                "total_uppercase": parsed.get("total_uppercase", ""),
                "total_lowercase": parsed.get("total_lowercase", ""),
                "total_verified": parsed.get("total_verified", False),
                "total_verify_msg": parsed.get("total_verify_msg", ""),
            }
            if info["total"] == 0 and info["amount"] > 0:
                info["total"] = round(info["amount"] + info["tax"], 2)
            invoices.append(info)

        # 创建Excel
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "发票明细"

        # 样式
        header_font = Font(bold=True, size=12, color="FFFFFF")
        header_fill = PatternFill(start_color="2563EB", end_color="2563EB", fill_type="solid")
        total_fill = PatternFill(start_color="F0F7FF", end_color="F0F7FF", fill_type="solid")
        thin_border = Border(
            left=Side(style='thin'), right=Side(style='thin'),
            top=Side(style='thin'), bottom=Side(style='thin')
        )
        center_align = Alignment(horizontal='center', vertical='center')

        # 标题行
        title = ws.cell(row=1, column=1, value="发票统计明细表")
        title.font = Font(bold=True, size=16)
        ws.merge_cells('A1:L1')

        # 汇总行
        ws.cell(row=2, column=1, value=f"共 {len(invoices)} 张发票")
        total_a = sum(inv["amount"] for inv in invoices)
        total_t = sum(inv["tax"] for inv in invoices)
        total_all = sum(inv["total"] for inv in invoices)
        ws.cell(row=2, column=4, value=f"金额合计：¥{total_a:.2f}")
        ws.cell(row=2, column=6, value=f"税额合计：¥{total_t:.2f}")
        ws.cell(row=2, column=8, value=f"价税合计：¥{total_all:.2f}")
        ws.cell(row=2, column=8).font = Font(bold=True, color="CC0000", size=12)

        # 表头
        headers = ["序号", "发票类型", "发票代码", "发票号码", "开票日期",
                   "销售方", "购买方", "金额", "税额", "价税合计",
                   "大写金额", "验证状态"]
        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=4, column=col, value=h)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = center_align
            cell.border = thin_border

        # 数据行
        for idx, inv in enumerate(invoices, 1):
            row = idx + 4
            verify_text = "✅ 验证通过" if inv["total_verified"] else inv.get("total_verify_msg", "")
            data = [idx, inv["invoice_type"], inv["invoice_code"], inv["invoice_number"],
                    inv["invoice_date"], inv["seller"], inv["buyer"],
                    inv["amount"], inv["tax"], inv["total"],
                    inv.get("total_uppercase", ""), verify_text]
            for col, val in enumerate(data, 1):
                cell = ws.cell(row=row, column=col, value=val)
                cell.border = thin_border
                if col >= 8 and col <= 10:
                    cell.number_format = '#,##0.00'
                    cell.alignment = Alignment(horizontal='right')
                elif col == 1:
                    cell.alignment = center_align
                elif col == 12:
                    # 验证状态列着色
                    if inv["total_verified"]:
                        cell.font = Font(color="2e7d32")
                    else:
                        cell.font = Font(color="e65100")

        # 合计行
        sum_row = len(invoices) + 5
        ws.cell(row=sum_row, column=1, value="合计").font = Font(bold=True)
        ws.cell(row=sum_row, column=8, value=total_a).font = Font(bold=True)
        ws.cell(row=sum_row, column=8).number_format = '#,##0.00'
        ws.cell(row=sum_row, column=9, value=total_t).font = Font(bold=True)
        ws.cell(row=sum_row, column=9).number_format = '#,##0.00'
        ws.cell(row=sum_row, column=10, value=total_all).font = Font(bold=True, color="CC0000")
        ws.cell(row=sum_row, column=10).number_format = '#,##0.00'
        for col in range(1, 13):
            ws.cell(row=sum_row, column=col).fill = total_fill
            ws.cell(row=sum_row, column=col).border = thin_border

        # 列宽
        col_widths = [6, 18, 14, 14, 12, 24, 24, 12, 12, 14, 22, 20]
        for i, w in enumerate(col_widths, 1):
            ws.column_dimensions[openpyxl.utils.get_column_letter(i)].width = w

        out_path = str(OUTPUT_DIR / f"invoice_stats_{gen_id()}.xlsx")
        wb.save(out_path)
        return send_file(out_path, as_attachment=True,
                       download_name="发票统计明细.xlsx")
    except Exception as e:
        return json_err(str(e))


# ============================================================
# 启动
# ============================================================

if __name__ == "__main__":
    HOST = os.environ.get("HOST", "0.0.0.0")
    PORT = int(os.environ.get("PORT", 5700))

    print("=" * 60)
    print("  PDFMaster v1.0 - 全能PDF处理软件")
    print("  功能：查看 | 编辑 | 注释 | 转换 | 压缩 | 合并拆分")
    print("        OCR | 水印 | 签名 | 加密 | 打印 | 页面管理 | 发票")
    print("=" * 60)

    try:
        from waitress import serve
        print("  服务器: Waitress (生产模式, 多线程)")
        print(f"  地址: http://{HOST}:{PORT}")
        print("=" * 60)
        serve(app, host=HOST, port=PORT, threads=8)
    except ImportError:
        print("  服务器: Flask 开发模式 (单线程)")
        print(f"  地址: http://{HOST}:{PORT}")
        print("=" * 60)
        app.run(host=HOST, port=PORT, debug=False)
